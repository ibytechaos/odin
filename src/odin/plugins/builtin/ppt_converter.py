"""PPT to HTML converter plugin for Odin.

This plugin converts PowerPoint presentations to beautiful HTML presentations
using LLM to redesign the content into reveal.js format.

Pure text mode - no vision/multimodal model required.
"""

import base64
from pathlib import Path
from typing import Annotated, Any

from openai import AsyncOpenAI
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import Field

from odin.config.settings import get_settings
from odin.decorators import tool
from odin.logging import get_logger
from odin.plugins import DecoratorPlugin, PluginConfig

logger = get_logger(__name__)


# System prompt for text-based PPT to HTML conversion
SYSTEM_PROMPT = """你是一位顶级的演示文稿设计师，擅长将内容转化为视觉震撼的 reveal.js 演示文稿。

## 你的设计风格

你的作品以以下特点著称：
1. **视觉层次分明** - 标题醒目，内容有主次之分
2. **色彩运用大胆** - 使用渐变、对比色创造视觉冲击
3. **排版呼吸感强** - 充足的留白，不拥挤
4. **动效优雅克制** - 适度的过渡动画，不喧宾夺主
5. **信息可视化** - 善用图标、列表、卡片等元素组织信息

## 必须遵守的技术规范

### HTML 结构
```html
<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>演示文稿标题</title>
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4/dist/reveal.css">
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4/dist/theme/black.css">
    <style>
        /* 在这里写自定义 CSS */
    </style>
</head>
<body>
    <div class="reveal">
        <div class="slides">
            <section>幻灯片内容</section>
        </div>
    </div>
    <script src="https://cdn.jsdelivr.net/npm/reveal.js@4/dist/reveal.js"></script>
    <script>
        Reveal.initialize({
            hash: true,
            transition: 'slide',
            backgroundTransition: 'fade'
        });
    </script>
</body>
</html>
```

### 自定义 CSS 示例（你应该创造更好的）

```css
:root {
    --primary-color: #667eea;
    --secondary-color: #764ba2;
    --text-color: #ffffff;
    --bg-dark: #1a1a2e;
}

.reveal {
    font-family: 'Segoe UI', 'PingFang SC', 'Microsoft YaHei', sans-serif;
}

/* 渐变背景 */
.reveal .slides section {
    background: linear-gradient(135deg, var(--bg-dark) 0%, #16213e 100%);
}

/* 标题样式 */
.reveal h1 {
    font-size: 2.5em;
    font-weight: 700;
    background: linear-gradient(90deg, var(--primary-color), var(--secondary-color));
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    text-shadow: none;
}

.reveal h2 {
    font-size: 1.8em;
    color: var(--primary-color);
    border-bottom: 3px solid var(--secondary-color);
    padding-bottom: 10px;
    display: inline-block;
}

/* 内容卡片 */
.card {
    background: rgba(255,255,255,0.1);
    border-radius: 16px;
    padding: 30px;
    backdrop-filter: blur(10px);
    border: 1px solid rgba(255,255,255,0.2);
    margin: 20px 0;
}

/* 列表美化 */
.reveal ul {
    list-style: none;
    padding-left: 0;
}

.reveal ul li {
    padding: 15px 20px;
    margin: 10px 0;
    background: rgba(255,255,255,0.05);
    border-left: 4px solid var(--primary-color);
    border-radius: 0 8px 8px 0;
}

/* 强调文字 */
.highlight {
    color: var(--secondary-color);
    font-weight: bold;
}

/* 数字/统计展示 */
.stat-number {
    font-size: 4em;
    font-weight: 800;
    background: linear-gradient(90deg, var(--primary-color), var(--secondary-color));
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
}

/* 两栏布局 */
.two-columns {
    display: grid;
    grid-template-columns: 1fr 1fr;
    gap: 40px;
    text-align: left;
}

/* 图片容器 */
.image-container {
    border-radius: 12px;
    overflow: hidden;
    box-shadow: 0 20px 40px rgba(0,0,0,0.3);
}

.image-container img {
    max-width: 100%;
    height: auto;
}
```

### 幻灯片类型设计指南

1. **封面页** - 大标题 + 副标题，可用全屏渐变背景
2. **目录页** - 清晰的章节列表，可带编号
3. **内容页** - 标题 + 要点列表或卡片
4. **图片页** - 图片居中，配简短说明
5. **对比页** - 两栏布局对比展示
6. **数据页** - 突出关键数字
7. **结束页** - 感谢语或联系方式

## 输出要求

1. 直接输出完整 HTML 代码，以 `<!DOCTYPE html>` 开头
2. 不要输出任何解释、说明或 markdown 代码块标记
3. 确保所有 CSS 都在 `<style>` 标签内
4. 图片使用提供的占位符格式：`<img src="{{IMAGE_SLIDE1_IMG1}}">`
5. 中文内容使用中文标点符号
"""


class PPTConverterPlugin(DecoratorPlugin):
    """PPT to HTML converter plugin.

    Converts PowerPoint presentations to beautiful HTML presentations
    using LLM to redesign the content. Pure text mode, no VLM required.
    """

    def __init__(self, config: PluginConfig | None = None) -> None:
        super().__init__(config)
        self._llm_client: AsyncOpenAI | None = None
        self._llm_model: str | None = None

    @property
    def name(self) -> str:
        return "ppt_converter"

    @property
    def version(self) -> str:
        return "1.0.0"

    @property
    def description(self) -> str:
        return "Convert PowerPoint to beautiful HTML presentations"

    async def initialize(self) -> None:
        """Initialize plugin with LLM client."""
        await super().initialize()
        settings = get_settings()

        self._llm_client = AsyncOpenAI(
            api_key=settings.openai_api_key,
            base_url=settings.openai_base_url,
        )
        self._llm_model = settings.openai_model

    def _extract_ppt_content(self, file_path: str) -> dict[str, Any]:
        """Extract all content from a PowerPoint file.

        Args:
            file_path: Path to the PPTX file

        Returns:
            Dictionary containing slide data
        """
        logger.info(f"[1/4] Opening PPT file: {file_path}")
        prs = Presentation(file_path)
        slides = []
        total_texts = 0
        total_images = 0

        logger.info(f"[1/4] PPT has {len(prs.slides)} slides, extracting content...")

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_data: dict[str, Any] = {
                "slide_number": slide_idx,
                "texts": [],
                "images": [],
                "notes": "",
            }

            # Extract text from all shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            is_title = shape == slide.shapes.title if hasattr(slide.shapes, "title") else False
                            slide_data["texts"].append({
                                "text": text,
                                "is_title": is_title,
                            })
                            total_texts += 1

                # Extract images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        content_type = image.content_type
                        b64_data = base64.b64encode(image_bytes).decode("utf-8")
                        data_url = f"data:{content_type};base64,{b64_data}"
                        slide_data["images"].append({
                            "data_url": data_url,
                            "content_type": content_type,
                        })
                        total_images += 1
                    except Exception as e:
                        logger.warning(f"Failed to extract image from slide {slide_idx}: {e}")

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_data["notes"] = notes_text

            slides.append(slide_data)

            # Log progress for each slide
            text_count = len(slide_data["texts"])
            img_count = len(slide_data["images"])
            logger.info(f"  Slide {slide_idx}: {text_count} texts, {img_count} images")

        title = prs.core_properties.title or "Untitled Presentation"

        logger.info(f"[1/4] Extraction complete: {len(slides)} slides, {total_texts} texts, {total_images} images")

        return {
            "title": title,
            "slides": slides,
            "slide_count": len(slides),
        }

    async def _generate_html(
        self,
        content: dict[str, Any],
        theme: str = "auto",
        language: str = "zh",
    ) -> str:
        """Generate HTML using extracted text content.

        Args:
            content: Extracted PPT content
            theme: Color theme preference
            language: Output language

        Returns:
            Complete HTML string
        """
        if not self._llm_client:
            raise RuntimeError("LLM client not initialized")

        # Format content as structured text for better LLM understanding
        lines = [
            f"# 演示文稿标题: {content['title']}",
            f"# 总共 {content['slide_count']} 页幻灯片",
            "",
            "---",
            ""
        ]

        image_placeholders: dict[str, str] = {}

        for slide in content["slides"]:
            lines.append(f"## 第 {slide['slide_number']} 页")
            lines.append("")

            # Extract text with better structure
            has_content = False
            if slide["texts"]:
                for text_item in slide["texts"]:
                    text = text_item['text']
                    if text_item.get("is_title"):
                        lines.append(f"**标题**: {text}")
                    else:
                        lines.append(f"- {text}")
                    has_content = True

            if not has_content and not slide["images"]:
                lines.append("(空白页)")

            # Handle images
            if slide["images"]:
                lines.append("")
                lines.append("**图片**:")
                for idx, img in enumerate(slide["images"]):
                    placeholder = f"{{{{IMAGE_SLIDE{slide['slide_number']}_IMG{idx + 1}}}}}"
                    image_placeholders[placeholder] = img["data_url"]
                    lines.append(f"  - 图片{idx + 1}: 使用占位符 {placeholder}")

            # Include notes if any
            if slide["notes"]:
                lines.append("")
                lines.append(f"**演讲者备注**: {slide['notes']}")

            lines.append("")
            lines.append("---")
            lines.append("")

        content_text = "\n".join(lines)

        # Build user prompt
        theme_instruction = {
            "dark": "使用深色主题（深色背景、浅色文字）",
            "light": "使用浅色主题（浅色背景、深色文字）",
            "auto": "根据内容自动选择合适的配色方案",
        }.get(theme, "根据内容自动选择合适的配色方案")

        user_prompt = f"""请将以下 PPT 内容转换为美观的 reveal.js HTML 演示文稿。

## PPT 原始内容

{content_text}

## 设计要求

1. **主题风格**: {theme_instruction}
2. **输出语言**: {language}
3. **设计风格**: 现代、专业、简洁

## 图片处理

如果 PPT 中有图片，请在 HTML 中使用以下格式：
<img src="{{{{IMAGE_SLIDE1_IMG1}}}}" alt="图片描述" style="max-width: 80%;">

系统会自动将占位符替换为实际的 base64 图片数据。

## 输出要求

直接输出完整的 HTML 代码，必须以 <!DOCTYPE html> 开头。
不要输出任何解释、说明或 markdown 代码块标记。"""

        logger.info("[3/4] Calling LLM API...")
        logger.info(f"  Model: {self._llm_model or 'gpt-4o'}")
        logger.info(f"  Prompt length: {len(user_prompt)} chars")
        logger.info(f"  Max tokens: 16000, Timeout: 180s")

        try:
            response = await self._llm_client.chat.completions.create(
                model=self._llm_model or "gpt-4o",
                messages=[
                    {"role": "system", "content": SYSTEM_PROMPT},
                    {"role": "user", "content": user_prompt},
                ],
                temperature=0.7,
                max_tokens=16000,
                timeout=180.0,
            )
            logger.info(f"[3/4] LLM API call successful")
        except Exception as e:
            logger.error(f"[3/4] LLM API call FAILED: {type(e).__name__}: {e}")
            raise

        html_content = response.choices[0].message.content or ""
        logger.info(f"[3/4] LLM response: {len(html_content)} chars")

        # Clean up response - remove markdown code block markers if present
        html_content = html_content.strip()
        if html_content.startswith("```html"):
            html_content = html_content[7:]
        if html_content.startswith("```"):
            html_content = html_content[3:]
        if html_content.endswith("```"):
            html_content = html_content[:-3]
        html_content = html_content.strip()

        # Replace image placeholders with actual base64 data URLs
        replaced_count = 0
        for placeholder, data_url in image_placeholders.items():
            if placeholder in html_content:
                html_content = html_content.replace(placeholder, data_url)
                replaced_count += 1

        if image_placeholders:
            logger.info(f"[3/4] Replaced {replaced_count}/{len(image_placeholders)} image placeholders")

        return html_content

    @tool(description="Convert PowerPoint to beautiful HTML presentation")
    async def ppt_to_html(
        self,
        file_path: Annotated[str, Field(description="Path to PPT/PPTX file")],
        output_path: Annotated[
            str | None,
            Field(description="Output HTML file path. If not provided, saves next to input file.")
        ] = None,
        theme: Annotated[
            str,
            Field(description="Color theme: dark, light, or auto (default: auto)")
        ] = "auto",
        language: Annotated[
            str,
            Field(description="Output language code (default: zh for Chinese)")
        ] = "zh",
    ) -> dict[str, Any]:
        """Convert a PowerPoint file to a beautiful HTML presentation.

        This tool:
        1. Extracts text and images from PPT
        2. Uses LLM to redesign content into reveal.js HTML
        3. Outputs a single self-contained HTML file

        No vision/multimodal model required - pure text mode.

        Args:
            file_path: Path to the PPTX file
            output_path: Where to save the HTML (optional)
            theme: Color theme preference
            language: Output language

        Returns:
            Success status, output path, and metadata
        """
        try:
            # Validate input file
            input_path = Path(file_path).expanduser().resolve()
            logger.info(f"=" * 50)
            logger.info(f"PPT to HTML Conversion Started")
            logger.info(f"=" * 50)
            logger.info(f"Input: {input_path}")

            if not input_path.exists():
                logger.error(f"File not found: {file_path}")
                return {"success": False, "error": f"File not found: {file_path}"}

            if input_path.suffix.lower() not in (".ppt", ".pptx"):
                logger.error(f"Invalid file type: {input_path.suffix}")
                return {"success": False, "error": f"Invalid file type: {input_path.suffix}"}

            # Determine output path
            out_path = Path(output_path).expanduser().resolve() if output_path else input_path.with_suffix(".html")
            logger.info(f"Output: {out_path}")
            logger.info(f"Theme: {theme}, Language: {language}")
            logger.info(f"-" * 50)

            # Extract content (text + embedded images)
            content = self._extract_ppt_content(str(input_path))

            # Initialize LLM if needed
            if not self._llm_client:
                logger.info(f"[2/4] Initializing LLM client...")
                await self.initialize()
                logger.info(f"[2/4] LLM client initialized")
            else:
                logger.info("[2/4] LLM client ready")

            # Generate HTML
            logger.info("[2/4] Preparing prompt for LLM...")
            html_content = await self._generate_html(
                content, theme=theme, language=language
            )

            # Validate HTML
            logger.info(f"[4/4] Processing LLM output...")
            if not html_content.strip().startswith("<!DOCTYPE") and not html_content.strip().startswith("<html"):
                logger.warning("[4/4] WARNING: LLM output may not be valid HTML")
                logger.warning(f"[4/4] Output starts with: {html_content[:100]}...")

            # Save HTML file
            out_path.parent.mkdir(parents=True, exist_ok=True)
            out_path.write_text(html_content, encoding="utf-8")
            logger.info(f"[4/4] HTML saved: {out_path}")
            logger.info(f"[4/4] File size: {out_path.stat().st_size} bytes")
            logger.info(f"=" * 50)
            logger.info(f"Conversion Complete!")
            logger.info(f"=" * 50)

            return {
                "success": True,
                "data": {
                    "input_file": str(input_path),
                    "output_file": str(out_path),
                    "slide_count": content["slide_count"],
                    "title": content["title"],
                    "theme": theme,
                    "file_size_bytes": out_path.stat().st_size,
                },
            }

        except Exception as e:
            logger.error(f"=" * 50)
            logger.error(f"PPT conversion FAILED: {type(e).__name__}: {e}")
            logger.error(f"=" * 50)
            return {"success": False, "error": str(e)}
