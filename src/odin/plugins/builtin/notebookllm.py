"""NotebookLLM Plugin for Odin.

Browser automation tools for interacting with Google NotebookLLM.
Provides capabilities to add notes, generate infographics and presentations,
and download generated content.

Uses browser_session for remote Chrome DevTools Protocol (CDP) connection.

Configuration via environment variables:
    CHROME_DEBUG_HOST: Remote Chrome host (e.g., chrome.example.com)
    CHROME_DEBUG_PORT: Remote Chrome port (default: 443 for TLS, 9222 otherwise)
    CHROME_DEBUG_TLS: Use TLS (true/false)

Requirements:
    - playwright
    - pdf2image (for PDF to images conversion)
    - poppler-utils (system dependency for pdf2image)

Tools:
- notebookllm_add_source: Add a web source (URL) to a notebook
- notebookllm_get_notebook_summary: Get summary info of a notebook (sources, artifacts)
- notebookllm_generate_mindmap: Generate a mind map from sources
- notebookllm_generate_infographic: Generate an infographic from sources
- notebookllm_generate_presentation: Generate a presentation from sources
- notebookllm_generate_and_download_infographic: Generate and download infographic in one operation
- notebookllm_generate_and_download_presentation: Generate and download presentation in one operation
- notebookllm_download_content: Download generated content
- pdf_to_images: Convert PDF to images
- images_to_editable_pptx: Convert slide images to editable PPTX
- notebookllm_close_browser: Close browser connection
"""


import asyncio
import base64
import time
from pathlib import Path
from typing import Annotated, Any

from pydantic import Field

from odin.decorators import tool
from odin.plugins import DecoratorPlugin, PluginConfig
from odin.utils.browser_session import (
    BrowserSession,
    cleanup_all_browser_sessions,
    get_browser_session,
)


class NotebookLLMPlugin(DecoratorPlugin):
    """NotebookLLM automation plugin for adding notes and generating content.

    This plugin provides browser automation tools for Google NotebookLLM,
    including adding sources, generating mind maps, infographics, and presentations.
    """

    def __init__(self, config: PluginConfig | None = None) -> None:
        super().__init__(config)
        self._session: BrowserSession | None = None

    @property
    def name(self) -> str:
        return "notebookllm"

    @property
    def version(self) -> str:
        return "1.0.0"

    @property
    def description(self) -> str:
        return "NotebookLLM automation tools for notes, infographics, and presentations"

    async def _get_session(self) -> BrowserSession:
        """Get browser session from pool (uses environment config)."""
        if self._session is None or self._session.page is None:
            self._session = await get_browser_session()
        return self._session

    async def _get_page(self) -> Any:
        """Get browser page from session."""
        session = await self._get_session()
        return session.page

    async def _close_browser(self) -> None:
        """Disconnect from browser (does not close the browser)."""
        await cleanup_all_browser_sessions()
        self._session = None

    async def _wait_for_notebookllm_ready(self, page: Any, timeout: int = 60) -> bool:
        """Wait for NotebookLLM page to be ready (logged in and loaded)."""
        try:
            await page.wait_for_selector(
                'section.source-panel, button.add-source-button, '
                'button[aria-label="添加来源"], button[aria-label="Add source"]',
                timeout=timeout * 1000,
            )
            return True
        except Exception:
            return False

    async def _create_new_notebook(self, notebook_name: str | None = None) -> str | None:
        """Create a new notebook and return its URL."""
        _ = notebook_name  # Reserved for future use
        page = await self._get_page()

        await page.goto("https://notebooklm.google.com/", wait_until="domcontentloaded")
        await asyncio.sleep(3)

        if "accounts.google.com" in page.url:
            return None

        create_btn = await page.query_selector(
            'button[aria-label="新建笔记本"], button[aria-label="New notebook"], '
            'button.create-new-button'
        )

        if not create_btn:
            create_btn = await page.query_selector('.create-new-action-button')

        if create_btn:
            await create_btn.click()
            await asyncio.sleep(4)

        current_url: str = page.url

        if "notebooklm.google.com/notebook/" in current_url:
            return current_url

        return None

    @tool(description="Add a web source (URL) to a NotebookLLM notebook")
    async def notebookllm_add_source(
        self,
        source_url: Annotated[str, Field(description="URL of the webpage to add as a source")],
        notebook_url: Annotated[
            str | None,
            Field(description="Full URL of the NotebookLLM notebook (optional, creates new if not provided)")
        ] = None,
        notebook_name: Annotated[
            str | None,
            Field(description="Name for new notebook (only used when creating new notebook)")
        ] = None,
        wait_for_processing: Annotated[
            bool,
            Field(description="Whether to wait for source processing to complete")
        ] = True,
        timeout: Annotated[
            int,
            Field(description="Maximum time to wait in seconds", ge=30, le=600)
        ] = 120,
    ) -> dict[str, Any]:
        """Add a web source (URL) to a NotebookLLM notebook.

        If notebook_url is not provided, creates a new notebook first.

        Args:
            source_url: URL of the webpage to add as a source
            notebook_url: Full URL of the NotebookLLM notebook (optional)
            notebook_name: Name for new notebook (only used when creating new notebook)
            wait_for_processing: Whether to wait for source processing to complete
            timeout: Maximum time to wait in seconds

        Returns:
            Status of the operation including notebook_url
        """
        try:
            page = await self._get_page()

            created_new_notebook = False
            final_notebook_url = notebook_url

            if not notebook_url:
                final_notebook_url = await self._create_new_notebook(notebook_name)

                if not final_notebook_url:
                    if "accounts.google.com" in page.url:
                        return {
                            "success": False,
                            "error": "Not logged in to Google. Please login manually in the browser window.",
                            "action_required": "login",
                            "notebook_url": None,
                        }
                    return {
                        "success": False,
                        "error": "Failed to create new notebook",
                        "notebook_url": None,
                    }

                created_new_notebook = True
            else:
                await page.goto(notebook_url, wait_until="domcontentloaded")
                await asyncio.sleep(2)

            if "accounts.google.com" in page.url:
                return {
                    "success": False,
                    "error": "Not logged in to Google. Please login manually in the browser window.",
                    "action_required": "login",
                    "notebook_url": final_notebook_url,
                }

            if not await self._wait_for_notebookllm_ready(page, timeout=30):
                return {
                    "success": False,
                    "error": "NotebookLLM page did not load properly",
                    "notebook_url": final_notebook_url,
                }

            existing_dialog = await page.query_selector('.upload-dialog-panel')

            if not existing_dialog:
                add_button = await page.query_selector(
                    'button.add-source-button, button[aria-label="添加来源"], '
                    'button[aria-label="Add source"]'
                )

                if not add_button:
                    return {
                        "success": False,
                        "error": "Could not find 'Add source' button",
                        "notebook_url": final_notebook_url,
                    }

                await add_button.click()
                await asyncio.sleep(1.5)

                existing_dialog = await page.wait_for_selector(
                    '.upload-dialog-panel',
                    timeout=10000,
                )

            if not existing_dialog:
                return {
                    "success": False,
                    "error": "Upload dialog did not appear",
                    "notebook_url": final_notebook_url,
                }

            url_input = await page.query_selector(
                '.upload-dialog-panel textarea.text-area, '
                '.cdk-overlay-pane textarea.text-area, '
                '.multi-urls-input-form textarea'
            )

            if not url_input:
                website_chip = await page.query_selector(
                    '.mat-mdc-chip:has(.mdc-evolution-chip__text-label:has-text("网站")), '
                    '.mat-mdc-chip:has(mat-icon:has-text("web"))'
                )

                if not website_chip:
                    website_chip = await page.query_selector(
                        '.chip-group__chip:has-text("网站"), '
                        '.chip-group__chip:has-text("Website")'
                    )

                if not website_chip:
                    return {
                        "success": False,
                        "error": "Could not find 'Website' option in add source dialog",
                        "notebook_url": final_notebook_url,
                    }

                await website_chip.click()
                await asyncio.sleep(1)

                url_input = await page.query_selector(
                    '.upload-dialog-panel textarea.text-area, '
                    '.cdk-overlay-pane textarea.text-area, '
                    '.multi-urls-input-form textarea'
                )

            if not url_input:
                return {
                    "success": False,
                    "error": "Could not find URL input field",
                    "notebook_url": final_notebook_url,
                }

            await url_input.fill(source_url)
            await asyncio.sleep(0.5)

            submit_button = await page.query_selector(
                '.upload-dialog-panel button:has-text("插入"), '
                '.cdk-overlay-pane button:has-text("插入"), '
                '.upload-dialog-panel button:has-text("Insert")'
            )

            if not submit_button:
                return {
                    "success": False,
                    "error": "Could not find 'Insert' button",
                    "notebook_url": final_notebook_url,
                }

            await asyncio.sleep(0.5)
            await submit_button.click()

            if wait_for_processing:
                start_time = time.time()

                while time.time() - start_time < timeout:
                    dialog_visible = await page.query_selector('.upload-dialog-panel')
                    if not dialog_visible:
                        await asyncio.sleep(2)
                        break

                    error = await page.query_selector(
                        '.mat-mdc-snack-bar-container:has-text("error"), '
                        '[role="alert"]:has-text("失败"), '
                        '[role="alert"]:has-text("failed")'
                    )
                    if error:
                        error_text = await error.inner_text()
                        return {
                            "success": False,
                            "error": f"NotebookLLM error: {error_text}",
                            "notebook_url": final_notebook_url,
                        }

                    await asyncio.sleep(1)

            final_notebook_url = page.url

            return {
                "success": True,
                "data": {
                    "message": f"Successfully added source: {source_url}",
                    "notebook_url": final_notebook_url,
                    "created_new_notebook": created_new_notebook,
                },
            }

        except ConnectionError as e:
            return {
                "success": False,
                "error": str(e),
                "hint": "Start Chrome with: google-chrome --remote-debugging-port=9222",
                "notebook_url": notebook_url,
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "notebook_url": notebook_url,
            }

    @tool(description="Get summary information of a NotebookLLM notebook")
    async def notebookllm_get_notebook_summary(
        self,
        notebook_url: Annotated[str, Field(description="Full URL of the NotebookLLM notebook")],
        timeout: Annotated[
            int,
            Field(description="Maximum time to wait in seconds", ge=10, le=120)
        ] = 30,
    ) -> dict[str, Any]:
        """Get summary information of a NotebookLLM notebook.

        Returns information about sources, generated artifacts (mind maps,
        infographics, presentations), and notebook metadata.

        Args:
            notebook_url: Full URL of the NotebookLLM notebook
            timeout: Maximum time to wait in seconds

        Returns:
            Notebook summary including sources and artifacts
        """
        try:
            page = await self._get_page()

            await page.goto(notebook_url, wait_until="domcontentloaded")
            await asyncio.sleep(3)

            if "accounts.google.com" in page.url:
                return {
                    "success": False,
                    "error": "Not logged in to Google. Please login manually in the browser window.",
                    "action_required": "login",
                    "notebook_url": notebook_url,
                }

            if not await self._wait_for_notebookllm_ready(page, timeout=timeout):
                return {
                    "success": False,
                    "error": "NotebookLLM page did not load properly",
                    "notebook_url": notebook_url,
                }

            # Extract notebook information using JavaScript
            notebook_info = await page.evaluate("""
                () => {
                    const result = {
                        title: '',
                        sources: [],
                        artifacts: {
                            mindmaps: 0,
                            infographics: 0,
                            presentations: 0,
                            audios: 0,
                            total: 0
                        },
                        notes: []
                    };

                    // Get notebook title
                    const titleEl = document.querySelector(
                        '.notebook-title, [class*="notebook-name"], h1'
                    );
                    if (titleEl) {
                        result.title = titleEl.textContent?.trim() || '';
                    }

                    // Get sources from source panel
                    const sourceItems = document.querySelectorAll(
                        '.source-panel .source-item, ' +
                        '[class*="source-list"] [class*="source-item"], ' +
                        '.sources-container [class*="item"]'
                    );
                    sourceItems.forEach(item => {
                        const titleEl = item.querySelector(
                            '[class*="title"], [class*="name"], .source-title'
                        );
                        const typeEl = item.querySelector(
                            '[class*="type"], mat-icon'
                        );
                        result.sources.push({
                            title: titleEl?.textContent?.trim() || 'Untitled',
                            type: typeEl?.textContent?.trim() || 'unknown'
                        });
                    });

                    // Count artifacts in studio panel
                    const studioPanel = document.querySelector('.studio-panel');
                    if (studioPanel) {
                        // Mind maps (flowchart icon)
                        const mindmaps = studioPanel.querySelectorAll(
                            'mat-icon.artifact-icon:has-text("flowchart"), ' +
                            '[class*="artifact"]:has(mat-icon:has-text("flowchart"))'
                        );
                        result.artifacts.mindmaps = mindmaps.length;

                        // Infographics (stacked_bar_chart icon)
                        const infographics = studioPanel.querySelectorAll(
                            'mat-icon.artifact-icon:has-text("stacked_bar_chart"), ' +
                            '[class*="artifact"]:has(mat-icon:has-text("stacked_bar_chart"))'
                        );
                        result.artifacts.infographics = infographics.length;

                        // Presentations (tablet icon)
                        const presentations = studioPanel.querySelectorAll(
                            'mat-icon.artifact-icon:has-text("tablet"), ' +
                            '[class*="artifact"]:has(mat-icon:has-text("tablet"))'
                        );
                        result.artifacts.presentations = presentations.length;

                        // Audio overviews (headphones icon)
                        const audios = studioPanel.querySelectorAll(
                            'mat-icon.artifact-icon:has-text("headphones"), ' +
                            '[class*="artifact"]:has(mat-icon:has-text("headphones"))'
                        );
                        result.artifacts.audios = audios.length;

                        result.artifacts.total = (
                            result.artifacts.mindmaps +
                            result.artifacts.infographics +
                            result.artifacts.presentations +
                            result.artifacts.audios
                        );
                    }

                    // Get notes/saved notes
                    const noteItems = document.querySelectorAll(
                        '.notes-panel .note-item, ' +
                        '[class*="note-list"] [class*="note-item"]'
                    );
                    noteItems.forEach(item => {
                        const content = item.querySelector(
                            '[class*="content"], [class*="text"]'
                        );
                        if (content) {
                            result.notes.push({
                                content: content.textContent?.trim().substring(0, 200) || ''
                            });
                        }
                    });

                    return result;
                }
            """)

            return {
                "success": True,
                "data": {
                    "notebook_url": notebook_url,
                    "title": notebook_info.get("title", ""),
                    "sources_count": len(notebook_info.get("sources", [])),
                    "sources": notebook_info.get("sources", []),
                    "artifacts": notebook_info.get("artifacts", {}),
                    "notes_count": len(notebook_info.get("notes", [])),
                    "notes": notebook_info.get("notes", []),
                },
            }

        except ConnectionError as e:
            return {
                "success": False,
                "error": str(e),
                "hint": "Start Chrome with: google-chrome --remote-debugging-port=9222",
                "notebook_url": notebook_url,
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "notebook_url": notebook_url,
            }

    @tool(description="Generate a mind map from NotebookLLM sources")
    async def notebookllm_generate_mindmap(
        self,
        notebook_url: Annotated[str, Field(description="Full URL of the NotebookLLM notebook")],
        timeout: Annotated[
            int,
            Field(description="Maximum time to wait for generation in seconds", ge=60, le=600)
        ] = 180,
    ) -> dict[str, Any]:
        """Generate a mind map from NotebookLLM sources.

        Args:
            notebook_url: Full URL of the NotebookLLM notebook
            timeout: Maximum time to wait for generation in seconds

        Returns:
            Status of the operation with mind map info
        """
        try:
            page = await self._get_page()

            await page.goto(notebook_url, wait_until="domcontentloaded")
            await asyncio.sleep(2)

            if not await self._wait_for_notebookllm_ready(page, timeout=30):
                return {
                    "success": False,
                    "error": "NotebookLLM page did not load properly",
                    "notebook_url": notebook_url,
                }

            mindmap_btn = await page.query_selector('button.mind-map-button')

            if not mindmap_btn:
                mindmap_btn = await page.query_selector(
                    'button:has-text("思维导图"), button:has-text("Mind map")'
                )

            if not mindmap_btn:
                mindmap_btn = await page.query_selector(
                    '.studio-panel button:has(mat-icon:has-text("flowchart"))'
                )

            if not mindmap_btn:
                return {
                    "success": False,
                    "error": "Could not find mind map generation button",
                    "notebook_url": notebook_url,
                }

            existing_artifacts = await page.query_selector_all(
                '.studio-panel mat-icon.artifact-icon:has-text("flowchart")'
            )
            initial_count = len(existing_artifacts)

            await mindmap_btn.click()
            await asyncio.sleep(2)

            start_time = time.time()

            while time.time() - start_time < timeout:
                studio_panel = await page.query_selector('.studio-panel')
                if studio_panel:
                    studio_text = await studio_panel.inner_text()

                    if "正在生成" in studio_text or "Generating" in studio_text:
                        await asyncio.sleep(3)
                        continue

                    current_artifacts = await page.query_selector_all(
                        '.studio-panel mat-icon.artifact-icon:has-text("flowchart")'
                    )
                    if len(current_artifacts) > initial_count:
                        return {
                            "success": True,
                            "data": {
                                "message": "Mind map generated successfully",
                                "notebook_url": notebook_url,
                            },
                        }

                error = await page.query_selector(
                    '.mat-mdc-snack-bar-container:has-text("error"), '
                    '[role="alert"]:has-text("失败"), '
                    '[role="alert"]:has-text("failed")'
                )
                if error:
                    error_text = await error.inner_text()
                    return {
                        "success": False,
                        "error": f"Generation failed: {error_text}",
                        "notebook_url": notebook_url,
                    }

                await asyncio.sleep(3)

            return {
                "success": False,
                "error": "Timeout waiting for mind map generation",
                "notebook_url": notebook_url,
            }

        except ConnectionError as e:
            return {
                "success": False,
                "error": str(e),
                "hint": "Start Chrome with: google-chrome --remote-debugging-port=9222",
                "notebook_url": notebook_url,
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "notebook_url": notebook_url,
            }

    @tool(description="Generate an infographic from NotebookLLM sources")
    async def notebookllm_generate_infographic(
        self,
        notebook_url: Annotated[str, Field(description="Full URL of the NotebookLLM notebook")],
        timeout: Annotated[
            int,
            Field(description="Maximum time to wait for generation in seconds", ge=60, le=600)
        ] = 180,
    ) -> dict[str, Any]:
        """Generate an infographic from NotebookLLM sources.

        Args:
            notebook_url: Full URL of the NotebookLLM notebook
            timeout: Maximum time to wait for generation in seconds

        Returns:
            Status of the operation with infographic info
        """
        try:
            page = await self._get_page()

            await page.goto(notebook_url, wait_until="domcontentloaded")
            await asyncio.sleep(2)

            if not await self._wait_for_notebookllm_ready(page, timeout=30):
                return {
                    "success": False,
                    "error": "NotebookLLM page did not load properly",
                    "notebook_url": notebook_url,
                }

            infographic_btn = await page.query_selector(
                '.create-artifact-button-container:has-text("信息图")'
            )

            if not infographic_btn:
                infographic_btn = await page.query_selector(
                    '.create-artifact-button-container:has(mat-icon:has-text("stacked_bar_chart"))'
                )

            if not infographic_btn:
                infographic_btn = await page.query_selector(
                    '.create-artifact-button-container:has-text("Infographic")'
                )

            if not infographic_btn:
                return {
                    "success": False,
                    "error": "Could not find infographic generation button",
                    "notebook_url": notebook_url,
                }

            existing_artifacts = await page.query_selector_all(
                '.studio-panel mat-icon.artifact-icon:has-text("stacked_bar_chart")'
            )
            initial_count = len(existing_artifacts)

            await infographic_btn.click()
            await asyncio.sleep(2)

            start_time = time.time()

            while time.time() - start_time < timeout:
                studio_panel = await page.query_selector('.studio-panel')
                if studio_panel:
                    studio_text = await studio_panel.inner_text()

                    if "正在生成" in studio_text or "Generating" in studio_text:
                        await asyncio.sleep(3)
                        continue

                    current_artifacts = await page.query_selector_all(
                        '.studio-panel mat-icon.artifact-icon:has-text("stacked_bar_chart")'
                    )
                    if len(current_artifacts) > initial_count:
                        return {
                            "success": True,
                            "data": {
                                "message": "Infographic generated successfully",
                                "notebook_url": notebook_url,
                            },
                        }

                error = await page.query_selector(
                    '.mat-mdc-snack-bar-container:has-text("error"), '
                    '[role="alert"]:has-text("失败")'
                )
                if error:
                    error_text = await error.inner_text()
                    return {
                        "success": False,
                        "error": f"Generation failed: {error_text}",
                        "notebook_url": notebook_url,
                    }

                await asyncio.sleep(3)

            return {
                "success": False,
                "error": "Timeout waiting for infographic generation",
                "notebook_url": notebook_url,
            }

        except ConnectionError as e:
            return {
                "success": False,
                "error": str(e),
                "hint": "Start Chrome with: google-chrome --remote-debugging-port=9222",
                "notebook_url": notebook_url,
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "notebook_url": notebook_url,
            }

    @tool(description="Generate and download an infographic from NotebookLLM sources in one operation")
    async def notebookllm_generate_and_download_infographic(
        self,
        notebook_url: Annotated[str, Field(description="Full URL of the NotebookLLM notebook")],
        output_dir: Annotated[
            str,
            Field(description="Directory to save downloaded file")
        ],
        timeout: Annotated[
            int,
            Field(description="Maximum time to wait for generation and download in seconds", ge=60, le=1800)
        ] = 600,
        poll_interval: Annotated[
            int,
            Field(description="Interval between download attempts in seconds", ge=5, le=60)
        ] = 10,
    ) -> dict[str, Any]:
        """Generate and download an infographic from NotebookLLM sources.

        This method combines generation and download into a single operation.
        It clicks the generate button and then polls for download availability.
        Success is determined by successful download (if you can download, generation is complete).

        Args:
            notebook_url: Full URL of the NotebookLLM notebook
            output_dir: Directory to save downloaded file
            timeout: Maximum time to wait in seconds
            poll_interval: Interval between download attempts in seconds

        Returns:
            Status with downloaded file path
        """
        try:
            page = await self._get_page()
            output_path = Path(output_dir)
            output_path.mkdir(parents=True, exist_ok=True)

            await page.goto(notebook_url, wait_until="domcontentloaded")
            await asyncio.sleep(2)

            if not await self._wait_for_notebookllm_ready(page, timeout=30):
                return {
                    "success": False,
                    "error": "NotebookLLM page did not load properly",
                    "notebook_url": notebook_url,
                }

            # Click generate button
            infographic_btn = await page.query_selector(
                '.create-artifact-button-container:has-text("信息图")'
            )
            if not infographic_btn:
                infographic_btn = await page.query_selector(
                    '.create-artifact-button-container:has(mat-icon:has-text("stacked_bar_chart"))'
                )
            if not infographic_btn:
                infographic_btn = await page.query_selector(
                    '.create-artifact-button-container:has-text("Infographic")'
                )

            if not infographic_btn:
                return {
                    "success": False,
                    "error": "Could not find infographic generation button",
                    "notebook_url": notebook_url,
                }

            await infographic_btn.click()
            await asyncio.sleep(3)

            # Poll for download availability
            start_time = time.time()

            while time.time() - start_time < timeout:
                # Check for errors
                error = await page.query_selector(
                    '.mat-mdc-snack-bar-container:has-text("error"), '
                    '[role="alert"]:has-text("失败"), '
                    '[role="alert"]:has-text("failed")'
                )
                if error:
                    error_text = await error.inner_text()
                    return {
                        "success": False,
                        "error": f"Generation failed: {error_text}",
                        "notebook_url": notebook_url,
                    }

                # Try to download
                result = await self._download_artifact(
                    page, "stacked_bar_chart", "infographic", output_path, 30
                )

                if result:
                    return {
                        "success": True,
                        "data": {
                            "message": "Infographic generated and downloaded successfully",
                            "notebook_url": notebook_url,
                            "file": result,
                            "output_dir": str(output_path),
                        },
                    }

                # Wait before next attempt
                await asyncio.sleep(poll_interval)

            return {
                "success": False,
                "error": "Timeout waiting for infographic generation and download",
                "notebook_url": notebook_url,
            }

        except ConnectionError as e:
            return {
                "success": False,
                "error": str(e),
                "hint": "Start Chrome with: google-chrome --remote-debugging-port=9222",
                "notebook_url": notebook_url,
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "notebook_url": notebook_url,
            }

    @tool(description="Generate a presentation (slides) from NotebookLLM sources")
    async def notebookllm_generate_presentation(
        self,
        notebook_url: Annotated[str, Field(description="Full URL of the NotebookLLM notebook")],
        timeout: Annotated[
            int,
            Field(description="Maximum time to wait for generation in seconds", ge=60, le=7200)
        ] = 3600,
    ) -> dict[str, Any]:
        """Generate a presentation (slides) from NotebookLLM sources.

        Note: Presentation generation can take a long time (up to 1 hour).

        Args:
            notebook_url: Full URL of the NotebookLLM notebook
            timeout: Maximum time to wait for generation in seconds (default: 1 hour)

        Returns:
            Status of the operation with presentation info
        """
        try:
            page = await self._get_page()

            await page.goto(notebook_url, wait_until="domcontentloaded")
            await asyncio.sleep(2)

            if not await self._wait_for_notebookllm_ready(page, timeout=30):
                return {
                    "success": False,
                    "error": "NotebookLLM page did not load properly",
                    "notebook_url": notebook_url,
                }

            presentation_btn = await page.query_selector(
                '.create-artifact-button-container:has-text("演示文稿")'
            )

            if not presentation_btn:
                presentation_btn = await page.query_selector(
                    '.create-artifact-button-container:has(mat-icon:has-text("tablet"))'
                )

            if not presentation_btn:
                presentation_btn = await page.query_selector(
                    '.create-artifact-button-container:has-text("Presentation")'
                )

            if not presentation_btn:
                return {
                    "success": False,
                    "error": "Could not find presentation generation button",
                    "notebook_url": notebook_url,
                }

            existing_artifacts = await page.query_selector_all(
                '.studio-panel mat-icon.artifact-icon:has-text("tablet")'
            )
            initial_count = len(existing_artifacts)

            await presentation_btn.click()
            await asyncio.sleep(2)

            start_time = time.time()

            while time.time() - start_time < timeout:
                studio_panel = await page.query_selector('.studio-panel')
                if studio_panel:
                    studio_text = await studio_panel.inner_text()

                    if "正在生成" in studio_text or "Generating" in studio_text:
                        await asyncio.sleep(3)
                        continue

                    current_artifacts = await page.query_selector_all(
                        '.studio-panel mat-icon.artifact-icon:has-text("tablet")'
                    )
                    if len(current_artifacts) > initial_count:
                        return {
                            "success": True,
                            "data": {
                                "message": "Presentation generated successfully",
                                "notebook_url": notebook_url,
                            },
                        }

                error = await page.query_selector(
                    '.mat-mdc-snack-bar-container:has-text("error"), '
                    '[role="alert"]:has-text("失败")'
                )
                if error:
                    error_text = await error.inner_text()
                    return {
                        "success": False,
                        "error": f"Generation failed: {error_text}",
                        "notebook_url": notebook_url,
                    }

                await asyncio.sleep(3)

            return {
                "success": False,
                "error": "Timeout waiting for presentation generation",
                "notebook_url": notebook_url,
            }

        except ConnectionError as e:
            return {
                "success": False,
                "error": str(e),
                "hint": "Start Chrome with: google-chrome --remote-debugging-port=9222",
                "notebook_url": notebook_url,
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "notebook_url": notebook_url,
            }

    @tool(description="Generate and download a presentation from NotebookLLM sources in one operation")
    async def notebookllm_generate_and_download_presentation(
        self,
        notebook_url: Annotated[str, Field(description="Full URL of the NotebookLLM notebook")],
        output_dir: Annotated[
            str,
            Field(description="Directory to save downloaded file")
        ],
        timeout: Annotated[
            int,
            Field(description="Maximum time to wait for generation and download in seconds", ge=60, le=7200)
        ] = 3600,
        poll_interval: Annotated[
            int,
            Field(description="Interval between download attempts in seconds", ge=5, le=120)
        ] = 30,
    ) -> dict[str, Any]:
        """Generate and download a presentation from NotebookLLM sources.

        This method combines generation and download into a single operation.
        It clicks the generate button and then polls for download availability.
        Success is determined by successful download (if you can download, generation is complete).

        Note: Presentation generation can take a long time (up to 1 hour).

        Args:
            notebook_url: Full URL of the NotebookLLM notebook
            output_dir: Directory to save downloaded file
            timeout: Maximum time to wait in seconds (default: 1 hour)
            poll_interval: Interval between download attempts in seconds

        Returns:
            Status with downloaded file path
        """
        try:
            page = await self._get_page()
            output_path = Path(output_dir)
            output_path.mkdir(parents=True, exist_ok=True)

            await page.goto(notebook_url, wait_until="domcontentloaded")
            await asyncio.sleep(2)

            if not await self._wait_for_notebookllm_ready(page, timeout=30):
                return {
                    "success": False,
                    "error": "NotebookLLM page did not load properly",
                    "notebook_url": notebook_url,
                }

            # Click generate button
            presentation_btn = await page.query_selector(
                '.create-artifact-button-container:has-text("演示文稿")'
            )
            if not presentation_btn:
                presentation_btn = await page.query_selector(
                    '.create-artifact-button-container:has(mat-icon:has-text("tablet"))'
                )
            if not presentation_btn:
                presentation_btn = await page.query_selector(
                    '.create-artifact-button-container:has-text("Presentation")'
                )

            if not presentation_btn:
                return {
                    "success": False,
                    "error": "Could not find presentation generation button",
                    "notebook_url": notebook_url,
                }

            await presentation_btn.click()
            await asyncio.sleep(3)

            # Poll for download availability
            start_time = time.time()

            while time.time() - start_time < timeout:
                # Check for errors
                error = await page.query_selector(
                    '.mat-mdc-snack-bar-container:has-text("error"), '
                    '[role="alert"]:has-text("失败"), '
                    '[role="alert"]:has-text("failed")'
                )
                if error:
                    error_text = await error.inner_text()
                    return {
                        "success": False,
                        "error": f"Generation failed: {error_text}",
                        "notebook_url": notebook_url,
                    }

                # Try to download
                result = await self._download_artifact(
                    page, "tablet", "presentation", output_path, 60
                )

                if result:
                    return {
                        "success": True,
                        "data": {
                            "message": "Presentation generated and downloaded successfully",
                            "notebook_url": notebook_url,
                            "file": result,
                            "output_dir": str(output_path),
                        },
                    }

                # Wait before next attempt
                await asyncio.sleep(poll_interval)

            return {
                "success": False,
                "error": "Timeout waiting for presentation generation and download",
                "notebook_url": notebook_url,
            }

        except ConnectionError as e:
            return {
                "success": False,
                "error": str(e),
                "hint": "Start Chrome with: google-chrome --remote-debugging-port=9222",
                "notebook_url": notebook_url,
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "notebook_url": notebook_url,
            }

    async def _wait_for_content_rendered(
        self,
        page: Any,
        content_type: str,
        timeout: int = 60,
    ) -> bool:
        """Wait for artifact content to be fully rendered before downloading."""
        start_time = time.time()

        while time.time() - start_time < timeout:
            loading_indicators = await page.query_selector_all(
                '.mat-progress-spinner, '
                '.mat-progress-bar, '
                '.loading-indicator, '
                '.spinner, '
                '[class*="loading"], '
                '[class*="spinner"]'
            )

            visible_loading = False
            for indicator in loading_indicators:
                is_visible = await indicator.is_visible()
                if is_visible:
                    visible_loading = True
                    break

            if visible_loading:
                await asyncio.sleep(1)
                continue

            if content_type == "presentation":
                slides = await page.query_selector_all(
                    '.slide-container, '
                    '.presentation-slide, '
                    '[class*="slide"], '
                    '.pdf-page, '
                    'canvas'
                )

                if slides:
                    for slide in slides:
                        box = await slide.bounding_box()
                        if box and box['width'] > 0 and box['height'] > 0:
                            await asyncio.sleep(3)
                            return True

            elif content_type == "infographic":
                content_elements = await page.query_selector_all(
                    'img[src], svg, canvas'
                )
                for elem in content_elements:
                    is_visible = await elem.is_visible()
                    if is_visible:
                        box = await elem.bounding_box()
                        if box and box['width'] > 100 and box['height'] > 100:
                            await asyncio.sleep(2)
                            return True

            await asyncio.sleep(2)

            loading_still = await page.query_selector(
                '.mat-progress-spinner:visible, '
                '.mat-progress-bar:visible'
            )
            if not loading_still:
                await asyncio.sleep(3)
                return True

        return False

    async def _download_artifact(
        self,
        page: Any,
        artifact_icon: str,
        content_type: str,
        output_path: Path,
        timeout: int,
    ) -> dict[str, Any] | None:
        """Download a specific artifact by clicking it and then the download button."""
        artifact_btn = await page.query_selector(
            f'button:has(mat-icon.artifact-icon:has-text("{artifact_icon}"))'
        )

        if not artifact_btn:
            return None

        await artifact_btn.click()
        await asyncio.sleep(2)

        render_timeout = min(timeout, 120)
        await self._wait_for_content_rendered(page, content_type, timeout=render_timeout)

        download_btn = await page.query_selector(
            'button[aria-label="下载"], '
            'button[aria-label="Download"], '
            'button:has(mat-icon:has-text("save_alt"))'
        )

        if not download_btn:
            close_btn = await page.query_selector('button:has(mat-icon:has-text("close"))')
            if close_btn:
                await close_btn.click()
                await asyncio.sleep(1)
            return None

        try:
            async with page.expect_download(timeout=timeout * 1000) as download_info:
                await download_btn.click()

            download = await download_info.value
            suggested_name = download.suggested_filename
            ext = Path(suggested_name).suffix if suggested_name else ".png"
            file_path = output_path / f"{content_type}_{int(time.time())}{ext}"
            await download.save_as(str(file_path))

            close_btn = await page.query_selector('button:has(mat-icon:has-text("close"))')
            if close_btn:
                await close_btn.click()
                await asyncio.sleep(1)

            return {
                "type": content_type,
                "path": str(file_path),
                "original_name": suggested_name,
            }
        except Exception:
            close_btn = await page.query_selector('button:has(mat-icon:has-text("close"))')
            if close_btn:
                await close_btn.click()
                await asyncio.sleep(1)
            return None

    @tool(description="Download generated content from NotebookLLM")
    async def notebookllm_download_content(
        self,
        notebook_url: Annotated[str, Field(description="Full URL of the NotebookLLM notebook")],
        content_type: Annotated[
            str,
            Field(description="Type of content to download: infographic, presentation, or all")
        ] = "all",
        output_dir: Annotated[
            str | None,
            Field(description="Directory to save downloaded files (default: current directory)")
        ] = None,
        timeout: Annotated[
            int,
            Field(description="Maximum time to wait for downloads in seconds", ge=30, le=600)
        ] = 120,
    ) -> dict[str, Any]:
        """Download generated content (infographic and/or presentation) from NotebookLLM.

        Args:
            notebook_url: Full URL of the NotebookLLM notebook
            content_type: Type of content to download ("infographic", "presentation", or "all")
            output_dir: Directory to save downloaded files (default: current directory)
            timeout: Maximum time to wait for downloads in seconds

        Returns:
            Paths to downloaded files and status
        """
        try:
            page = await self._get_page()

            output_path = Path(output_dir) if output_dir else Path.cwd()
            output_path.mkdir(parents=True, exist_ok=True)

            await page.goto(notebook_url, wait_until="domcontentloaded")
            await asyncio.sleep(2)

            if not await self._wait_for_notebookllm_ready(page, timeout=30):
                return {
                    "success": False,
                    "error": "NotebookLLM page did not load properly",
                }

            downloaded_files = []

            if content_type in ("infographic", "all"):
                result = await self._download_artifact(
                    page, "stacked_bar_chart", "infographic", output_path, timeout
                )
                if result:
                    downloaded_files.append(result)

            if content_type in ("presentation", "all"):
                result = await self._download_artifact(
                    page, "tablet", "presentation", output_path, timeout
                )
                if result:
                    downloaded_files.append(result)

            if not downloaded_files:
                return {
                    "success": False,
                    "error": "No downloadable content found. Make sure infographic/presentation has been generated.",
                }

            return {
                "success": True,
                "data": {
                    "message": f"Downloaded {len(downloaded_files)} file(s)",
                    "files": downloaded_files,
                    "output_dir": str(output_path),
                },
            }

        except ConnectionError as e:
            return {
                "success": False,
                "error": str(e),
                "hint": "Start Chrome with: google-chrome --remote-debugging-port=9222",
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
            }

    @tool(description="Convert PDF file to a list of images (one per page)")
    async def pdf_to_images(
        self,
        pdf_path: Annotated[str, Field(description="Path to the PDF file")],
        output_dir: Annotated[
            str | None,
            Field(description="Directory to save images (default: same directory as PDF)")
        ] = None,
        dpi: Annotated[int, Field(description="Image resolution in DPI", ge=72, le=600)] = 200,
        format: Annotated[str, Field(description="Output image format (png, jpg, jpeg)")] = "png",
        return_base64: Annotated[
            bool,
            Field(description="Whether to include base64 encoded images in response")
        ] = False,
    ) -> dict[str, Any]:
        """Convert PDF file to a list of images (one per page).

        Args:
            pdf_path: Path to the PDF file
            output_dir: Directory to save images (default: same directory as PDF)
            dpi: Image resolution in DPI
            format: Output image format (png, jpg, jpeg)
            return_base64: Whether to include base64 encoded images in response

        Returns:
            List of generated image paths and optionally base64 data
        """
        try:
            from pdf2image import convert_from_path
        except ImportError:
            return {
                "success": False,
                "error": "pdf2image library not installed. Install with: pip install pdf2image",
                "hint": "Also requires poppler-utils: brew install poppler (macOS) or apt-get install poppler-utils (Linux)",
            }

        pdf_file = Path(pdf_path)
        if not pdf_file.exists():
            return {
                "success": False,
                "error": f"PDF file not found: {pdf_path}",
            }

        output_path = Path(output_dir) if output_dir else pdf_file.parent
        output_path.mkdir(parents=True, exist_ok=True)

        try:
            images = convert_from_path(str(pdf_file), dpi=dpi)

            image_files = []
            base64_images = []

            for i, image in enumerate(images):
                filename = f"{pdf_file.stem}_page_{i + 1}.{format}"
                image_path = output_path / filename
                image.save(str(image_path), format.upper())

                image_info = {
                    "page": i + 1,
                    "path": str(image_path),
                    "width": image.width,
                    "height": image.height,
                }
                image_files.append(image_info)

                if return_base64:
                    import io

                    buffer = io.BytesIO()
                    image.save(buffer, format=format.upper())
                    b64 = base64.b64encode(buffer.getvalue()).decode("utf-8")
                    base64_images.append({
                        "page": i + 1,
                        "base64": b64,
                        "mime_type": f"image/{format}",
                    })

            result_data: dict[str, Any] = {
                "message": f"Converted {len(images)} pages to images",
                "total_pages": len(images),
                "images": image_files,
                "output_dir": str(output_path),
                "format": format,
                "dpi": dpi,
            }

            if return_base64:
                result_data["base64_images"] = base64_images

            return {
                "success": True,
                "data": result_data,
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"PDF conversion failed: {e!s}",
            }

    def _detect_layout_with_yolo(
        self,
        image_path: str,
        conf_threshold: float = 0.25,
    ) -> list[dict[str, Any]]:
        """Detect layout elements using DocLayout-YOLO model.

        Returns a list of detected elements with type, bbox, and confidence.
        Types: title, text, figure, table, list, caption, etc.
        """
        try:
            from doclayout_yolo import YOLOv10
            from huggingface_hub import hf_hub_download
        except ImportError:
            raise ImportError(
                "doclayout-yolo not installed. Install with: pip install doclayout-yolo huggingface-hub"
            ) from None

        model_id = "juliozhao/DocLayout-YOLO-DocStructBench"
        model_file = "doclayout_yolo_docstructbench_imgsz1024.pt"

        model_path = hf_hub_download(repo_id=model_id, filename=model_file)
        model = YOLOv10(model_path)

        results = model.predict(image_path, imgsz=1024, conf=conf_threshold, verbose=False)

        elements = []
        for result in results:
            boxes = result.boxes
            names = result.names

            for i, box in enumerate(boxes):
                xyxy = box.xyxy[0].cpu().numpy()
                cls_id = int(box.cls[0])
                conf = float(box.conf[0])
                cls_name = names.get(cls_id, "text").lower()

                elements.append({
                    "id": f"elem_{i}",
                    "type": cls_name,
                    "bbox": {
                        "x": int(xyxy[0]),
                        "y": int(xyxy[1]),
                        "w": int(xyxy[2] - xyxy[0]),
                        "h": int(xyxy[3] - xyxy[1]),
                    },
                    "confidence": conf,
                })

        elements.sort(key=lambda e: (e["bbox"]["y"], e["bbox"]["x"]))
        return elements

    def _detect_background_color(self, image_path: str) -> str:
        """Detect dominant background color from image corners."""
        import numpy as np
        from PIL import Image

        img = Image.open(image_path)
        img_array = np.array(img)
        h, w = img_array.shape[:2]

        margin = min(50, h // 10, w // 10)
        samples = []

        samples.append(img_array[:margin, :margin])
        samples.append(img_array[:margin, -margin:])
        samples.append(img_array[-margin:, :margin])
        samples.append(img_array[-margin:, -margin:])

        all_pixels = np.concatenate([s.reshape(-1, 3) for s in samples])
        median_color = np.median(all_pixels, axis=0).astype(int)

        return f"#{median_color[0]:02x}{median_color[1]:02x}{median_color[2]:02x}"

    def _detect_text_color_from_region(
        self,
        image_path: str,
        bbox: dict[str, int],
    ) -> str:
        """Detect text color from a region (darkest significant color)."""
        import numpy as np
        from PIL import Image

        img = Image.open(image_path)
        img_array = np.array(img)

        x, y, w, h = bbox["x"], bbox["y"], bbox["w"], bbox["h"]
        region = img_array[y:y+h, x:x+w]

        if region.size == 0:
            return "#000000"

        pixels = region.reshape(-1, 3)
        luminance = 0.299 * pixels[:, 0] + 0.587 * pixels[:, 1] + 0.114 * pixels[:, 2]

        threshold = np.percentile(luminance, 15)
        dark_pixels = pixels[luminance <= threshold]

        if len(dark_pixels) == 0:
            return "#000000"

        median_color = np.median(dark_pixels, axis=0).astype(int)
        return f"#{median_color[0]:02x}{median_color[1]:02x}{median_color[2]:02x}"

    def _hex_to_rgb(self, hex_color: str) -> tuple[int, int, int]:
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip("#")
        return (
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )

    def _extract_figure_image(
        self,
        image_path: str,
        bbox: dict[str, int],
        output_path: str,
    ) -> str:
        """Extract figure region as a separate image."""
        from PIL import Image

        img = Image.open(image_path)
        x, y, w, h = bbox["x"], bbox["y"], bbox["w"], bbox["h"]
        cropped = img.crop((x, y, x + w, y + h))
        cropped.save(output_path)
        return output_path

    def _extract_text_with_ocr(
        self,
        image_path: str,
        bbox: dict[str, int],
        lang: str = "ch",
    ) -> str:
        """Extract text from a region using EasyOCR."""
        try:
            import easyocr
        except ImportError:
            return ""

        try:
            import tempfile

            from PIL import Image

            img = Image.open(image_path)
            x, y, w, h = bbox["x"], bbox["y"], bbox["w"], bbox["h"]
            cropped = img.crop((x, y, x + w, y + h))

            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
                cropped.save(f.name)
                temp_path = f.name

            # Map lang code to EasyOCR languages
            langs = ["en"]
            if lang == "ch":
                langs = ["ch_sim", "en"]
            elif lang == "cht":
                langs = ["ch_tra", "en"]

            reader = easyocr.Reader(langs, gpu=False)
            result = reader.readtext(temp_path)

            Path(temp_path).unlink(missing_ok=True)

            # EasyOCR returns list of (box, text, confidence)
            if result:
                texts = [item[1] for item in result if item and len(item) >= 2]
                return " ".join(texts)

            return ""
        except Exception:
            return ""

    @tool(description="Convert slide images to an editable PowerPoint presentation")
    async def images_to_editable_pptx(
        self,
        image_paths: Annotated[list[str], Field(description="List of paths to slide images (in order)")],
        output_path: Annotated[str, Field(description="Path for the output .pptx file")],
        analysis_data: Annotated[
            list[dict[str, Any]] | None,
            Field(description="Optional pre-analyzed layout data from multimodal LLM (list of slide analysis)")
        ] = None,
        use_ocr: Annotated[
            bool,
            Field(description="Use PaddleOCR to extract text when analysis_data is not provided")
        ] = True,
        ocr_lang: Annotated[
            str,
            Field(description="OCR language ('ch' for Chinese+English, 'en' for English only)")
        ] = "ch",
        slide_width_inches: Annotated[
            float,
            Field(description="Slide width in inches (default: 16:9 widescreen)", ge=1.0, le=30.0)
        ] = 13.333,
        min_font_size: Annotated[
            int,
            Field(description="Minimum font size in points", ge=6, le=72)
        ] = 10,
        max_font_size: Annotated[
            int,
            Field(description="Maximum font size in points", ge=12, le=144)
        ] = 72,
        default_font: Annotated[
            str,
            Field(description="Default font family for text")
        ] = "Microsoft YaHei",
        conf_threshold: Annotated[
            float,
            Field(description="Confidence threshold for layout detection", ge=0.1, le=1.0)
        ] = 0.25,
    ) -> dict[str, Any]:
        """Convert slide images to an editable PowerPoint presentation.

        This tool uses DocLayout-YOLO for layout detection to:
        1. Detect layout elements (title, text, figure, table, etc.)
        2. Keep figures/graphics as images (100% fidelity)
        3. Extract text using PaddleOCR (or use provided analysis_data)
        4. Create editable text boxes with detected styles

        For best results, provide pre-analyzed layout data from a multimodal LLM
        that includes text content and styling information.

        Args:
            image_paths: List of paths to slide images (in order)
            output_path: Path for the output .pptx file
            analysis_data: Optional pre-analyzed layout data from multimodal LLM
            use_ocr: Use PaddleOCR to extract text (default: True)
            ocr_lang: OCR language ('ch' for Chinese+English, 'en' for English only)
            slide_width_inches: Slide width in inches (default: 16:9 widescreen)
            min_font_size: Minimum font size in points
            max_font_size: Maximum font size in points
            default_font: Default font family for text
            conf_threshold: Confidence threshold for layout detection

        Returns:
            Status with output path and processing details
        """
        try:
            from PIL import Image
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt
        except ImportError as e:
            missing = str(e).split("'")[1] if "'" in str(e) else str(e)
            return {
                "success": False,
                "error": f"Missing dependency: {missing}",
                "hint": "Install with: pip install python-pptx pillow",
            }

        output_file = Path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]

        processed_slides: list[dict[str, Any]] = []
        temp_files: list[Path] = []

        try:
            for idx, image_path in enumerate(image_paths):
                slide_info: dict[str, Any] = {
                    "page": idx + 1,
                    "image": image_path,
                    "elements": 0,
                    "figures": 0,
                    "text_boxes": 0,
                    "status": "processing",
                }

                if not Path(image_path).exists():
                    slide_info["status"] = "error"
                    slide_info["error"] = "Image file not found"
                    processed_slides.append(slide_info)
                    continue

                img = Image.open(image_path)
                img_width, img_height = img.size

                aspect = img_width / img_height
                slide_height_inches = slide_width_inches / aspect

                if idx == 0:
                    prs.slide_width = Inches(slide_width_inches)
                    prs.slide_height = Inches(slide_height_inches)

                scale = slide_width_inches / img_width

                slide_analysis = None
                if analysis_data and idx < len(analysis_data):
                    slide_analysis = analysis_data[idx]

                if slide_analysis:
                    elements = slide_analysis.get("elements", [])
                    bg_color = slide_analysis.get("background_color", "#FFFFFF")
                else:
                    try:
                        elements = self._detect_layout_with_yolo(image_path, conf_threshold)
                    except ImportError:
                        elements = []
                    except Exception as e:
                        slide_info["status"] = "warning"
                        slide_info["warning"] = f"Layout detection failed: {e!s}"
                        elements = []
                    bg_color = self._detect_background_color(image_path)

                slide_info["elements"] = len(elements)

                slide = prs.slides.add_slide(blank_layout)

                bg = slide.background
                fill = bg.fill
                fill.solid()
                r, g, b = self._hex_to_rgb(bg_color)
                fill.fore_color.rgb = RGBColor(r, g, b)

                figure_types = {"figure", "table", "chart", "image", "picture"}

                for elem in elements:
                    bbox = elem.get("bbox", {})
                    elem_type = elem.get("type", "text").lower()

                    if not bbox:
                        continue

                    x = bbox.get("x", 0)
                    y = bbox.get("y", 0)
                    w = bbox.get("w", bbox.get("width", 100))
                    h = bbox.get("h", bbox.get("height", 50))

                    left = Inches(x * scale)
                    top = Inches(y * scale)
                    width = Inches(w * scale)
                    height = Inches(h * scale)

                    if elem_type in figure_types or elem.get("is_figure"):
                        temp_fig_path = output_file.parent / f"_temp_fig_{idx}_{elem.get('id', 'fig')}.png"
                        self._extract_figure_image(image_path, {"x": x, "y": y, "w": w, "h": h}, str(temp_fig_path))
                        temp_files.append(temp_fig_path)

                        slide.shapes.add_picture(str(temp_fig_path), left, top, width, height)
                        slide_info["figures"] += 1
                    else:
                        text_segments = elem.get("text_segments", [])

                        if not text_segments:
                            text_content = elem.get("text_content", elem.get("text", ""))

                            # Use OCR to extract text if not provided
                            if not text_content and use_ocr:
                                text_content = self._extract_text_with_ocr(
                                    image_path, {"x": x, "y": y, "w": w, "h": h}, ocr_lang
                                )

                            if text_content:
                                text_color = self._detect_text_color_from_region(
                                    image_path, {"x": x, "y": y, "w": w, "h": h}
                                )
                                font_size_pt = min(max_font_size, max(min_font_size, int(h * scale * 72 * 0.6)))
                                is_bold = elem_type in ("title", "header")

                                text_segments = [{
                                    "text": text_content,
                                    "style": {
                                        "font_family": default_font,
                                        "font_size_pt": font_size_pt,
                                        "color_hex": text_color,
                                        "is_bold": is_bold,
                                        "alignment": "center" if elem_type == "title" else "left",
                                    }
                                }]

                        if text_segments:
                            textbox = slide.shapes.add_textbox(left, top, width, height)
                            frame = textbox.text_frame
                            frame.word_wrap = False

                            para = frame.paragraphs[0]

                            first_style = text_segments[0].get("style", {})
                            align = first_style.get("alignment", "left")
                            if align == "center":
                                para.alignment = PP_ALIGN.CENTER
                            elif align == "right":
                                para.alignment = PP_ALIGN.RIGHT
                            else:
                                para.alignment = PP_ALIGN.LEFT

                            for seg in text_segments:
                                run = para.add_run()
                                run.text = seg.get("text", "")

                                style = seg.get("style", {})
                                font_size = style.get("font_size_pt", 12)
                                font_size = min(max_font_size, max(min_font_size, font_size))
                                run.font.size = Pt(font_size)
                                run.font.bold = style.get("is_bold", False)
                                run.font.italic = style.get("is_italic", False)
                                run.font.name = style.get("font_family", default_font)

                                color_hex = style.get("color_hex", "#000000")
                                r, g, b = self._hex_to_rgb(color_hex)
                                run.font.color.rgb = RGBColor(r, g, b)

                            slide_info["text_boxes"] += 1

                slide_info["status"] = "success"
                processed_slides.append(slide_info)

            prs.save(str(output_file))

            for temp_file in temp_files:
                temp_file.unlink(missing_ok=True)

            success_count = sum(1 for s in processed_slides if s["status"] == "success")
            total_elements = sum(s.get("elements", 0) for s in processed_slides)
            total_figures = sum(s.get("figures", 0) for s in processed_slides)
            total_text_boxes = sum(s.get("text_boxes", 0) for s in processed_slides)

            return {
                "success": True,
                "data": {
                    "message": f"Created editable PPTX with {success_count}/{len(image_paths)} slides",
                    "output_path": str(output_file),
                    "total_slides": len(image_paths),
                    "successful_slides": success_count,
                    "total_elements": total_elements,
                    "total_figures": total_figures,
                    "total_text_boxes": total_text_boxes,
                    "slides": processed_slides,
                },
            }

        except Exception as e:
            for temp_file in temp_files:
                temp_file.unlink(missing_ok=True)

            return {
                "success": False,
                "error": f"PPTX generation failed: {e!s}",
            }

    @tool(description="Close the browser connection")
    async def notebookllm_close_browser(self) -> dict[str, Any]:
        """Close the browser connection (does not close the actual browser).

        Returns:
            Status of the operation
        """
        try:
            await self._close_browser()
            return {
                "success": True,
                "data": {"message": "Browser connection closed successfully"},
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
            }
