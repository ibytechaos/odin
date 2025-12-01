"""NotebookLLM Automation Tools for Odin.

Browser automation tools for interacting with Google NotebookLLM.
Provides capabilities to add notes, generate infographics and presentations,
and download generated content.

Uses the shared BrowserManager utility for Chrome DevTools Protocol (CDP) connection.
Start Chrome with: google-chrome --remote-debugging-port=9222

Requirements:
    - playwright
    - pdf2image (for PDF to images conversion)
    - poppler-utils (system dependency for pdf2image)
"""

import asyncio
import base64
import time
from pathlib import Path
from typing import Any

from odin.decorators import tool
from odin.plugins import DecoratorPlugin
from odin.utils.browser import get_browser_manager


class NotebookLLMTools(DecoratorPlugin):
    """NotebookLLM automation tools for adding notes and generating content."""

    def __init__(self):
        super().__init__()
        self._browser_manager = get_browser_manager()

    @property
    def name(self) -> str:
        return "notebookllm"

    @property
    def version(self) -> str:
        return "1.0.0"

    @property
    def description(self) -> str:
        return "NotebookLLM automation tools for notes, infographics, and presentations"

    async def _get_page(self):
        """Get browser page from shared BrowserManager."""
        return await self._browser_manager.get_page()

    async def _close_browser(self) -> None:
        """Disconnect from browser (does not close the browser)."""
        await self._browser_manager.close()

    async def _wait_for_notebookllm_ready(self, page, timeout: int = 60) -> bool:
        """Wait for NotebookLLM page to be ready (logged in and loaded)."""
        try:
            # Wait for main notebook page elements to load
            # Based on actual page inspection: source-panel section or add-source button
            await page.wait_for_selector(
                'section.source-panel, button.add-source-button, '
                'button[aria-label="添加来源"], button[aria-label="Add source"]',
                timeout=timeout * 1000,
            )
            return True
        except Exception:
            return False

    async def _create_new_notebook(self, notebook_name: str | None = None) -> str | None:
        """Create a new notebook and return its URL.

        Args:
            notebook_name: Optional name for the notebook (not used - NotebookLLM creates with default name)

        Returns:
            URL of the newly created notebook, or None if failed
        """
        page = await self._get_page()

        # Navigate to NotebookLLM home page
        # Use domcontentloaded instead of networkidle to avoid timeout
        await page.goto("https://notebooklm.google.com/", wait_until="domcontentloaded")
        await asyncio.sleep(3)  # Give time for dynamic content to load

        # Check if logged in
        if "accounts.google.com" in page.url:
            return None

        # Find and click "New notebook" button
        # Based on actual page: button[aria-label="新建笔记本"] or button.create-new-button
        create_btn = await page.query_selector(
            'button[aria-label="新建笔记本"], button[aria-label="New notebook"], '
            'button.create-new-button'
        )

        if not create_btn:
            # Try the card-style create button
            create_btn = await page.query_selector('.create-new-action-button')

        if create_btn:
            await create_btn.click()
            # Wait for navigation to new notebook page
            await asyncio.sleep(4)

        # Get current URL
        current_url = page.url

        # Check if we're now on a notebook page
        if "notebooklm.google.com/notebook/" in current_url:
            return current_url

        return None

    @tool()
    async def notebookllm_add_source(
        self,
        source_url: str,
        notebook_url: str | None = None,
        notebook_name: str | None = None,
        wait_for_processing: bool = True,
        timeout: int = 120,
    ) -> dict[str, Any]:
        """Add a web source (URL) to a NotebookLLM notebook.

        If notebook_url is not provided, creates a new notebook first.

        Args:
            source_url: URL of the webpage to add as a source
            notebook_url: Full URL of the NotebookLLM notebook (optional, creates new if not provided)
            notebook_name: Name for new notebook (only used when creating new notebook)
            wait_for_processing: Whether to wait for source processing to complete
            timeout: Maximum time to wait in seconds

        Returns:
            Status of the operation including notebook_url (always returned)
        """
        try:
            page = await self._get_page()

            created_new_notebook = False
            final_notebook_url = notebook_url

            # If no notebook URL provided, create a new one
            if not notebook_url:
                final_notebook_url = await self._create_new_notebook(notebook_name)

                if not final_notebook_url:
                    # Check if login is needed
                    if "accounts.google.com" in page.url:
                        return {
                            "success": False,
                            "error": "Not logged in to Google. Please login manually in the browser window.",
                            "action_required": "login",
                            "notebook_url": None,
                        }
                    return {
                        "success": False,
                        "error": "Failed to create new notebook",
                        "notebook_url": None,
                    }

                created_new_notebook = True
            else:
                # Navigate to existing notebook
                await page.goto(notebook_url, wait_until="domcontentloaded")
                await asyncio.sleep(2)  # Give time for dynamic content to load

            # Check if logged in
            if "accounts.google.com" in page.url:
                return {
                    "success": False,
                    "error": "Not logged in to Google. Please login manually in the browser window.",
                    "action_required": "login",
                    "notebook_url": final_notebook_url,
                }

            # Wait for page to load
            if not await self._wait_for_notebookllm_ready(page, timeout=30):
                return {
                    "success": False,
                    "error": "NotebookLLM page did not load properly",
                    "notebook_url": final_notebook_url,
                }

            # Step 1: Check if add source dialog is already open (common for new notebooks)
            # If it's already open, we can skip clicking the add source button
            existing_dialog = await page.query_selector('.upload-dialog-panel')

            if not existing_dialog:
                # Dialog not open, click "Add source" button to open it
                add_button = await page.query_selector(
                    'button.add-source-button, button[aria-label="添加来源"], '
                    'button[aria-label="Add source"]'
                )

                if not add_button:
                    return {
                        "success": False,
                        "error": "Could not find 'Add source' button",
                        "notebook_url": final_notebook_url,
                    }

                await add_button.click()
                await asyncio.sleep(1.5)

                # Wait for upload dialog to appear
                existing_dialog = await page.wait_for_selector(
                    '.upload-dialog-panel',
                    timeout=10000,
                )

            if not existing_dialog:
                return {
                    "success": False,
                    "error": "Upload dialog did not appear",
                    "notebook_url": final_notebook_url,
                }

            # Step 2: Check if we're already in the URL input view (has textarea)
            # If yes, skip clicking the website chip
            url_input = await page.query_selector(
                '.upload-dialog-panel textarea.text-area, '
                '.cdk-overlay-pane textarea.text-area, '
                '.multi-urls-input-form textarea'
            )

            if not url_input:
                # Not in URL input view, need to click the "网站" (Website) chip
                website_chip = await page.query_selector(
                    '.mat-mdc-chip:has(.mdc-evolution-chip__text-label:has-text("网站")), '
                    '.mat-mdc-chip:has(mat-icon:has-text("web"))'
                )

                if not website_chip:
                    # Try alternative selector
                    website_chip = await page.query_selector(
                        '.chip-group__chip:has-text("网站"), '
                        '.chip-group__chip:has-text("Website")'
                    )

                if not website_chip:
                    return {
                        "success": False,
                        "error": "Could not find 'Website' option in add source dialog",
                        "notebook_url": final_notebook_url,
                    }

                await website_chip.click()
                await asyncio.sleep(1)

                # Now find the URL textarea
                url_input = await page.query_selector(
                    '.upload-dialog-panel textarea.text-area, '
                    '.cdk-overlay-pane textarea.text-area, '
                    '.multi-urls-input-form textarea'
                )

            # Step 3: Fill the URL textarea
            if not url_input:
                return {
                    "success": False,
                    "error": "Could not find URL input field",
                    "notebook_url": final_notebook_url,
                }

            # Enter the source URL
            await url_input.fill(source_url)
            await asyncio.sleep(0.5)

            # Step 4: Click the "插入" (Insert) submit button
            # IMPORTANT: Must select the button INSIDE the dialog, not the one outside
            # There are multiple submit-button elements on the page
            submit_button = await page.query_selector(
                '.upload-dialog-panel button:has-text("插入"), '
                '.cdk-overlay-pane button:has-text("插入"), '
                '.upload-dialog-panel button:has-text("Insert")'
            )

            if not submit_button:
                return {
                    "success": False,
                    "error": "Could not find 'Insert' button",
                    "notebook_url": final_notebook_url,
                }

            # Wait for button to be enabled (it's disabled when input is empty)
            await asyncio.sleep(0.5)
            await submit_button.click()

            # Step 5: Wait for processing if requested
            if wait_for_processing:
                # Wait for dialog to close (indicates source was submitted)
                start_time = time.time()

                while time.time() - start_time < timeout:
                    # Check if dialog closed
                    dialog_visible = await page.query_selector('.upload-dialog-panel')
                    if not dialog_visible:
                        # Dialog closed, source is being processed
                        # Wait a moment for the page to update
                        await asyncio.sleep(2)
                        break

                    # Check for error snackbar or alert
                    error = await page.query_selector(
                        '.mat-mdc-snack-bar-container:has-text("error"), '
                        '[role="alert"]:has-text("失败"), '
                        '[role="alert"]:has-text("failed")'
                    )
                    if error:
                        error_text = await error.inner_text()
                        return {
                            "success": False,
                            "error": f"NotebookLLM error: {error_text}",
                            "notebook_url": final_notebook_url,
                        }

                    await asyncio.sleep(1)

            # Get the current URL
            final_notebook_url = page.url

            return {
                "success": True,
                "message": f"Successfully added source: {source_url}",
                "notebook_url": final_notebook_url,
                "created_new_notebook": created_new_notebook,
            }

        except ConnectionError as e:
            return {
                "success": False,
                "error": str(e),
                "hint": "Start Chrome with: google-chrome --remote-debugging-port=9222",
                "notebook_url": notebook_url,
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "notebook_url": notebook_url,
            }

    @tool()
    async def notebookllm_generate_mindmap(
        self,
        notebook_url: str,
        timeout: int = 180,
    ) -> dict[str, Any]:
        """Generate a mind map from NotebookLLM sources.

        The mind map is generated in the Studio panel based on all selected sources.
        Sources are selected by default when you add them to the notebook.

        Args:
            notebook_url: Full URL of the NotebookLLM notebook
            timeout: Maximum time to wait for generation in seconds

        Returns:
            Status of the operation with mind map info
        """
        try:
            page = await self._get_page()

            # Navigate to notebook
            await page.goto(notebook_url, wait_until="domcontentloaded")
            await asyncio.sleep(2)  # Wait for page to stabilize

            if not await self._wait_for_notebookllm_ready(page, timeout=30):
                return {
                    "success": False,
                    "error": "NotebookLLM page did not load properly",
                    "notebook_url": notebook_url,
                }

            # Step 1: Find and click the mind map button in Studio panel
            # Based on actual page: button.mind-map-button with flowchart icon
            mindmap_btn = await page.query_selector('button.mind-map-button')

            if not mindmap_btn:
                # Try by text
                mindmap_btn = await page.query_selector(
                    'button:has-text("思维导图"), button:has-text("Mind map")'
                )

            if not mindmap_btn:
                # Try by icon in studio panel
                mindmap_btn = await page.query_selector(
                    '.studio-panel button:has(mat-icon:has-text("flowchart"))'
                )

            if not mindmap_btn:
                return {
                    "success": False,
                    "error": "Could not find mind map generation button. Make sure the Studio panel is visible.",
                    "notebook_url": notebook_url,
                }

            # Count existing mindmap artifacts before clicking
            # Artifacts have mat-icon with artifact-icon class and flowchart text
            existing_artifacts = await page.query_selector_all(
                '.studio-panel mat-icon.artifact-icon:has-text("flowchart")'
            )
            initial_count = len(existing_artifacts)

            await mindmap_btn.click()
            await asyncio.sleep(2)

            # Step 2: Wait for generation to complete
            start_time = time.time()

            while time.time() - start_time < timeout:
                # Check Studio panel for completion indicators
                studio_panel = await page.query_selector('.studio-panel')
                if studio_panel:
                    studio_text = await studio_panel.inner_text()

                    # Check if still generating - just wait
                    if "正在生成" in studio_text or "Generating" in studio_text:
                        await asyncio.sleep(3)
                        continue

                    # Check for new artifact (might have completed quickly)
                    current_artifacts = await page.query_selector_all(
                        '.studio-panel mat-icon.artifact-icon:has-text("flowchart")'
                    )
                    if len(current_artifacts) > initial_count:
                        return {
                            "success": True,
                            "message": "Mind map generated successfully",
                            "notebook_url": notebook_url,
                        }

                # Check for error messages
                error = await page.query_selector(
                    '.mat-mdc-snack-bar-container:has-text("error"), '
                    '[role="alert"]:has-text("失败"), '
                    '[role="alert"]:has-text("failed")'
                )
                if error:
                    error_text = await error.inner_text()
                    return {
                        "success": False,
                        "error": f"Generation failed: {error_text}",
                        "notebook_url": notebook_url,
                    }

                await asyncio.sleep(3)

            return {
                "success": False,
                "error": "Timeout waiting for mind map generation",
                "notebook_url": notebook_url,
            }

        except ConnectionError as e:
            return {
                "success": False,
                "error": str(e),
                "hint": "Start Chrome with: google-chrome --remote-debugging-port=9222",
                "notebook_url": notebook_url,
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "notebook_url": notebook_url,
            }

    @tool()
    async def notebookllm_generate_infographic(
        self,
        notebook_url: str,
        timeout: int = 180,
    ) -> dict[str, Any]:
        """Generate an infographic from NotebookLLM sources.

        The infographic is generated in the Studio panel based on all selected sources.

        Args:
            notebook_url: Full URL of the NotebookLLM notebook
            timeout: Maximum time to wait for generation in seconds

        Returns:
            Status of the operation with infographic info
        """
        try:
            page = await self._get_page()

            # Navigate to notebook
            await page.goto(notebook_url, wait_until="domcontentloaded")
            await asyncio.sleep(2)

            if not await self._wait_for_notebookllm_ready(page, timeout=30):
                return {
                    "success": False,
                    "error": "NotebookLLM page did not load properly",
                    "notebook_url": notebook_url,
                }

            # Step 1: Find and click the infographic button in Studio panel
            # The buttons are divs with class create-artifact-button-container
            infographic_btn = await page.query_selector(
                '.create-artifact-button-container:has-text("信息图")'
            )

            if not infographic_btn:
                # Try by icon
                infographic_btn = await page.query_selector(
                    '.create-artifact-button-container:has(mat-icon:has-text("stacked_bar_chart"))'
                )

            if not infographic_btn:
                # Try English text
                infographic_btn = await page.query_selector(
                    '.create-artifact-button-container:has-text("Infographic")'
                )

            if not infographic_btn:
                return {
                    "success": False,
                    "error": "Could not find infographic generation button. Make sure the Studio panel is visible.",
                    "notebook_url": notebook_url,
                }

            # Count existing infographic artifacts before clicking
            # Artifacts have mat-icon with artifact-icon class
            existing_artifacts = await page.query_selector_all(
                '.studio-panel mat-icon.artifact-icon:has-text("stacked_bar_chart")'
            )
            initial_count = len(existing_artifacts)

            await infographic_btn.click()
            await asyncio.sleep(2)

            # Step 2: Wait for generation to complete
            start_time = time.time()

            while time.time() - start_time < timeout:
                studio_panel = await page.query_selector('.studio-panel')
                if studio_panel:
                    studio_text = await studio_panel.inner_text()

                    # Check if still generating - just wait
                    if "正在生成" in studio_text or "Generating" in studio_text:
                        await asyncio.sleep(3)
                        continue

                    # Check for new artifact (might have completed quickly)
                    current_artifacts = await page.query_selector_all(
                        '.studio-panel mat-icon.artifact-icon:has-text("stacked_bar_chart")'
                    )
                    if len(current_artifacts) > initial_count:
                        return {
                            "success": True,
                            "message": "Infographic generated successfully",
                            "notebook_url": notebook_url,
                        }

                # Check for error
                error = await page.query_selector(
                    '.mat-mdc-snack-bar-container:has-text("error"), '
                    '[role="alert"]:has-text("失败")'
                )
                if error:
                    error_text = await error.inner_text()
                    return {
                        "success": False,
                        "error": f"Generation failed: {error_text}",
                        "notebook_url": notebook_url,
                    }

                await asyncio.sleep(3)

            return {
                "success": False,
                "error": "Timeout waiting for infographic generation",
                "notebook_url": notebook_url,
            }

        except ConnectionError as e:
            return {
                "success": False,
                "error": str(e),
                "hint": "Start Chrome with: google-chrome --remote-debugging-port=9222",
                "notebook_url": notebook_url,
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "notebook_url": notebook_url,
            }

    @tool()
    async def notebookllm_generate_presentation(
        self,
        notebook_url: str,
        timeout: int = 3600,
    ) -> dict[str, Any]:
        """Generate a presentation (slides) from NotebookLLM sources.

        The presentation is generated in the Studio panel based on all selected sources.
        Note: Presentation generation can take a long time (up to 1 hour).

        Args:
            notebook_url: Full URL of the NotebookLLM notebook
            timeout: Maximum time to wait for generation in seconds (default: 1 hour)

        Returns:
            Status of the operation with presentation info
        """
        try:
            page = await self._get_page()

            # Navigate to notebook
            await page.goto(notebook_url, wait_until="domcontentloaded")
            await asyncio.sleep(2)

            if not await self._wait_for_notebookllm_ready(page, timeout=30):
                return {
                    "success": False,
                    "error": "NotebookLLM page did not load properly",
                    "notebook_url": notebook_url,
                }

            # Step 1: Find and click the presentation button in Studio panel
            # The buttons are divs with class create-artifact-button-container
            presentation_btn = await page.query_selector(
                '.create-artifact-button-container:has-text("演示文稿")'
            )

            if not presentation_btn:
                # Try by icon
                presentation_btn = await page.query_selector(
                    '.create-artifact-button-container:has(mat-icon:has-text("tablet"))'
                )

            if not presentation_btn:
                # Try English text
                presentation_btn = await page.query_selector(
                    '.create-artifact-button-container:has-text("Presentation")'
                )

            if not presentation_btn:
                return {
                    "success": False,
                    "error": "Could not find presentation generation button. Make sure the Studio panel is visible.",
                    "notebook_url": notebook_url,
                }

            # Count existing presentation artifacts before clicking
            # Artifacts have mat-icon with artifact-icon class
            existing_artifacts = await page.query_selector_all(
                '.studio-panel mat-icon.artifact-icon:has-text("tablet")'
            )
            initial_count = len(existing_artifacts)

            await presentation_btn.click()
            await asyncio.sleep(2)

            # Step 2: Wait for generation to complete
            start_time = time.time()

            while time.time() - start_time < timeout:
                studio_panel = await page.query_selector('.studio-panel')
                if studio_panel:
                    studio_text = await studio_panel.inner_text()

                    # Check if still generating - just wait
                    if "正在生成" in studio_text or "Generating" in studio_text:
                        await asyncio.sleep(3)
                        continue

                    # Check for new artifact (might have completed quickly)
                    current_artifacts = await page.query_selector_all(
                        '.studio-panel mat-icon.artifact-icon:has-text("tablet")'
                    )
                    if len(current_artifacts) > initial_count:
                        return {
                            "success": True,
                            "message": "Presentation generated successfully",
                            "notebook_url": notebook_url,
                        }

                # Check for error
                error = await page.query_selector(
                    '.mat-mdc-snack-bar-container:has-text("error"), '
                    '[role="alert"]:has-text("失败")'
                )
                if error:
                    error_text = await error.inner_text()
                    return {
                        "success": False,
                        "error": f"Generation failed: {error_text}",
                        "notebook_url": notebook_url,
                    }

                await asyncio.sleep(3)

            return {
                "success": False,
                "error": "Timeout waiting for presentation generation",
                "notebook_url": notebook_url,
            }

        except ConnectionError as e:
            return {
                "success": False,
                "error": str(e),
                "hint": "Start Chrome with: google-chrome --remote-debugging-port=9222",
                "notebook_url": notebook_url,
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "notebook_url": notebook_url,
            }

    async def _wait_for_content_rendered(
        self,
        page,
        content_type: str,
        timeout: int = 60,
    ) -> bool:
        """Wait for artifact content to be fully rendered before downloading.

        For presentations/infographics, the content may take time to render,
        especially for large documents. This method waits until:
        1. Loading indicators disappear
        2. Content elements are visible
        3. No ongoing network activity related to content loading

        Args:
            page: Playwright page object
            content_type: Type of content ("presentation" or "infographic")
            timeout: Maximum time to wait in seconds

        Returns:
            True if content appears ready, False if timeout
        """
        start_time = time.time()

        while time.time() - start_time < timeout:
            # Check for loading spinners/indicators
            loading_indicators = await page.query_selector_all(
                '.mat-progress-spinner, '
                '.mat-progress-bar, '
                '.loading-indicator, '
                '.spinner, '
                '[class*="loading"], '
                '[class*="spinner"]'
            )

            # Filter to only visible loading indicators
            visible_loading = False
            for indicator in loading_indicators:
                is_visible = await indicator.is_visible()
                if is_visible:
                    visible_loading = True
                    break

            if visible_loading:
                await asyncio.sleep(1)
                continue

            # For presentations, check if slides are rendered
            if content_type == "presentation":
                # Look for slide content containers
                slides = await page.query_selector_all(
                    '.slide-container, '
                    '.presentation-slide, '
                    '[class*="slide"], '
                    '.pdf-page, '
                    'canvas'
                )

                if slides:
                    # Additional check: wait for at least one slide to have content
                    # by checking if canvas has non-zero dimensions or content is visible
                    for slide in slides:
                        box = await slide.bounding_box()
                        if box and box['width'] > 0 and box['height'] > 0:
                            # Found a rendered slide, wait a bit more for all to load
                            await asyncio.sleep(3)
                            return True

            # For infographics, check if the image/svg is loaded
            elif content_type == "infographic":
                content_elements = await page.query_selector_all(
                    'img[src], svg, canvas'
                )
                for elem in content_elements:
                    is_visible = await elem.is_visible()
                    if is_visible:
                        box = await elem.bounding_box()
                        if box and box['width'] > 100 and box['height'] > 100:
                            await asyncio.sleep(2)
                            return True

            # General check: wait for network to be idle
            # If no specific content found, wait a bit and check again
            await asyncio.sleep(2)

            # After waiting, if no loading indicators, assume ready
            loading_still = await page.query_selector(
                '.mat-progress-spinner:visible, '
                '.mat-progress-bar:visible'
            )
            if not loading_still:
                # Extra safety wait for large content
                await asyncio.sleep(3)
                return True

        return False

    async def _download_artifact(
        self,
        page,
        artifact_icon: str,
        content_type: str,
        output_path: Path,
        timeout: int,
    ) -> dict[str, Any] | None:
        """Download a specific artifact by clicking it and then the download button.

        Args:
            page: Playwright page object
            artifact_icon: The mat-icon text for the artifact (e.g., "stacked_bar_chart")
            content_type: Type name for logging (e.g., "infographic")
            output_path: Directory to save the file
            timeout: Download timeout in seconds

        Returns:
            Download info dict or None if not found
        """
        # Find and click the artifact button in studio panel
        artifact_btn = await page.query_selector(
            f'button:has(mat-icon.artifact-icon:has-text("{artifact_icon}"))'
        )

        if not artifact_btn:
            return None

        # Click to open the artifact detail view
        await artifact_btn.click()
        await asyncio.sleep(2)

        # Wait for content to be fully rendered before downloading
        # This is critical for large presentations that take time to render
        render_timeout = min(timeout, 120)  # Cap render wait at 2 minutes
        content_ready = await self._wait_for_content_rendered(
            page, content_type, timeout=render_timeout
        )

        if not content_ready:
            # Log warning but still try to download
            # Sometimes content detection may fail but content is actually ready
            pass

        # Now find the download button (aria-label="下载" or "Download")
        download_btn = await page.query_selector(
            'button[aria-label="下载"], '
            'button[aria-label="Download"], '
            'button:has(mat-icon:has-text("save_alt"))'
        )

        if not download_btn:
            # Try to close the detail view before returning
            close_btn = await page.query_selector('button:has(mat-icon:has-text("close"))')
            if close_btn:
                await close_btn.click()
                await asyncio.sleep(1)
            return None

        # Click download and wait for the file
        try:
            async with page.expect_download(timeout=timeout * 1000) as download_info:
                await download_btn.click()

            download = await download_info.value
            suggested_name = download.suggested_filename
            # Determine extension from suggested filename or default
            ext = Path(suggested_name).suffix if suggested_name else ".png"
            file_path = output_path / f"{content_type}_{int(time.time())}{ext}"
            await download.save_as(str(file_path))

            # Close the detail view
            close_btn = await page.query_selector('button:has(mat-icon:has-text("close"))')
            if close_btn:
                await close_btn.click()
                await asyncio.sleep(1)

            return {
                "type": content_type,
                "path": str(file_path),
                "original_name": suggested_name,
            }
        except Exception:
            # Close the detail view on error
            close_btn = await page.query_selector('button:has(mat-icon:has-text("close"))')
            if close_btn:
                await close_btn.click()
                await asyncio.sleep(1)
            return None

    @tool()
    async def notebookllm_download_content(
        self,
        notebook_url: str,
        content_type: str = "all",
        output_dir: str | None = None,
        timeout: int = 120,
    ) -> dict[str, Any]:
        """Download generated content (infographic and/or presentation) from NotebookLLM.

        This tool clicks on the artifact in the Studio panel to open its detail view,
        then clicks the download button to save the file.

        Args:
            notebook_url: Full URL of the NotebookLLM notebook
            content_type: Type of content to download ("infographic", "presentation", or "all")
            output_dir: Directory to save downloaded files (default: current directory)
            timeout: Maximum time to wait for downloads in seconds

        Returns:
            Paths to downloaded files and status
        """
        try:
            page = await self._get_page()

            output_path = Path(output_dir) if output_dir else Path.cwd()
            output_path.mkdir(parents=True, exist_ok=True)

            # Navigate to notebook
            await page.goto(notebook_url, wait_until="domcontentloaded")
            await asyncio.sleep(2)

            if not await self._wait_for_notebookllm_ready(page, timeout=30):
                return {
                    "success": False,
                    "error": "NotebookLLM page did not load properly",
                }

            downloaded_files = []

            # Download infographic (icon: stacked_bar_chart)
            if content_type in ("infographic", "all"):
                result = await self._download_artifact(
                    page, "stacked_bar_chart", "infographic", output_path, timeout
                )
                if result:
                    downloaded_files.append(result)

            # Download presentation (icon: tablet)
            if content_type in ("presentation", "all"):
                result = await self._download_artifact(
                    page, "tablet", "presentation", output_path, timeout
                )
                if result:
                    downloaded_files.append(result)

            if not downloaded_files:
                return {
                    "success": False,
                    "error": "No downloadable content found. Make sure infographic/presentation has been generated.",
                }

            return {
                "success": True,
                "message": f"Downloaded {len(downloaded_files)} file(s)",
                "files": downloaded_files,
                "output_dir": str(output_path),
            }

        except ConnectionError as e:
            return {
                "success": False,
                "error": str(e),
                "hint": "Start Chrome with: google-chrome --remote-debugging-port=9222",
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
            }

    @tool()
    async def pdf_to_images(
        self,
        pdf_path: str,
        output_dir: str | None = None,
        dpi: int = 200,
        format: str = "png",
        return_base64: bool = False,
    ) -> dict[str, Any]:
        """Convert PDF file to a list of images (one per page).

        Args:
            pdf_path: Path to the PDF file
            output_dir: Directory to save images (default: same directory as PDF)
            dpi: Image resolution in DPI
            format: Output image format (png, jpg, jpeg)
            return_base64: Whether to include base64 encoded images in response

        Returns:
            List of generated image paths and optionally base64 data
        """
        try:
            from pdf2image import convert_from_path
        except ImportError:
            return {
                "success": False,
                "error": "pdf2image library not installed. Install with: pip install pdf2image",
                "hint": "Also requires poppler-utils: brew install poppler (macOS) or apt-get install poppler-utils (Linux)",
            }

        pdf_file = Path(pdf_path)
        if not pdf_file.exists():
            return {
                "success": False,
                "error": f"PDF file not found: {pdf_path}",
            }

        output_path = Path(output_dir) if output_dir else pdf_file.parent
        output_path.mkdir(parents=True, exist_ok=True)

        try:
            # Convert PDF to images
            images = convert_from_path(str(pdf_file), dpi=dpi)

            image_files = []
            base64_images = []

            for i, image in enumerate(images):
                # Save image
                filename = f"{pdf_file.stem}_page_{i + 1}.{format}"
                image_path = output_path / filename
                image.save(str(image_path), format.upper())

                image_info = {
                    "page": i + 1,
                    "path": str(image_path),
                    "width": image.width,
                    "height": image.height,
                }
                image_files.append(image_info)

                # Generate base64 if requested
                if return_base64:
                    import io

                    buffer = io.BytesIO()
                    image.save(buffer, format=format.upper())
                    b64 = base64.b64encode(buffer.getvalue()).decode("utf-8")
                    base64_images.append({
                        "page": i + 1,
                        "base64": b64,
                        "mime_type": f"image/{format}",
                    })

            result = {
                "success": True,
                "message": f"Converted {len(images)} pages to images",
                "total_pages": len(images),
                "images": image_files,
                "output_dir": str(output_path),
                "format": format,
                "dpi": dpi,
            }

            if return_base64:
                result["base64_images"] = base64_images

            return result

        except Exception as e:
            return {
                "success": False,
                "error": f"PDF conversion failed: {str(e)}",
            }

    def _detect_text_regions(
        self,
        image_path: str,
        lang: str = "ch",
    ) -> list[dict[str, Any]]:
        """Detect text regions in an image using PaddleOCR.

        Args:
            image_path: Path to the image file
            lang: Language code ('ch' for Chinese+English, 'en' for English only)

        Returns:
            List of detected text regions with coordinates, text, and confidence
        """
        try:
            # Apply compatibility patches for paddleocr 3.x (requires langchain shim)
            from odin.compat import patch_langchain_for_paddlex
            patch_langchain_for_paddlex()
            from paddleocr import PaddleOCR
        except ImportError:
            raise ImportError(
                "PaddleOCR not installed. Install with: pip install paddlepaddle paddleocr"
            ) from None

        # Initialize PaddleOCR
        # Note: PaddleOCR 3.x API changed - uses different params than 2.x
        ocr = PaddleOCR(lang=lang)

        # Run OCR - PaddleOCR 3.x returns results directly
        result = ocr.ocr(image_path)

        text_regions = []
        if result and result[0]:
            for line in result[0]:
                # line format: [box_coordinates, (text, confidence)]
                box = line[0]  # [[x1,y1], [x2,y2], [x3,y3], [x4,y4]]
                text, confidence = line[1]

                # Calculate bounding box (min/max of polygon points)
                x_coords = [p[0] for p in box]
                y_coords = [p[1] for p in box]
                x_min, x_max = min(x_coords), max(x_coords)
                y_min, y_max = min(y_coords), max(y_coords)

                # Estimate font size based on box height
                box_height = y_max - y_min
                # Rough estimation: font size ≈ box height * 0.75
                estimated_font_size = int(box_height * 0.75)

                text_regions.append({
                    "text": text,
                    "confidence": confidence,
                    "box": {
                        "x": int(x_min),
                        "y": int(y_min),
                        "width": int(x_max - x_min),
                        "height": int(y_max - y_min),
                    },
                    "polygon": [[int(p[0]), int(p[1])] for p in box],
                    "estimated_font_size": estimated_font_size,
                })

        return text_regions

    def _create_text_mask(
        self,
        image_size: tuple[int, int],
        text_regions: list[dict[str, Any]],
        dilate_pixels: int = 5,
    ):
        """Create a binary mask for text regions.

        Args:
            image_size: (width, height) of the image
            text_regions: List of text regions from OCR
            dilate_pixels: Pixels to dilate the mask (helps with inpainting)

        Returns:
            numpy array mask (255 for text areas, 0 for background)
        """
        import cv2
        import numpy as np

        width, height = image_size
        mask = np.zeros((height, width), dtype=np.uint8)

        for region in text_regions:
            # Use polygon for more accurate masking
            polygon = np.array(region["polygon"], dtype=np.int32)
            cv2.fillPoly(mask, [polygon], 255)

        # Dilate mask to ensure complete text coverage
        if dilate_pixels > 0:
            kernel = np.ones((dilate_pixels, dilate_pixels), np.uint8)
            mask = cv2.dilate(mask, kernel, iterations=1)

        return mask

    def _inpaint_background(
        self,
        image_path: str,
        mask,
        method: str = "telea",
        inpaint_radius: int = 5,
    ):
        """Remove text from image using inpainting to reconstruct background.

        Args:
            image_path: Path to the original image
            mask: Binary mask (255 for areas to inpaint)
            method: Inpainting method ('telea', 'ns', or 'lama')
            inpaint_radius: Radius for OpenCV inpainting methods

        Returns:
            Inpainted image (numpy array in BGR format)
        """
        import cv2

        image = cv2.imread(image_path)

        if method == "lama":
            # Try to use LaMa for better quality (requires lama-cleaner)
            try:
                from lama_cleaner.model_manager import ModelManager
                from lama_cleaner.schema import Config

                model = ModelManager(name="lama", device="cpu")
                config = Config(
                    ldm_steps=25,
                    hd_strategy="Original",
                    hd_strategy_crop_margin=128,
                    hd_strategy_crop_trigger_size=800,
                    hd_strategy_resize_limit=800,
                )
                # Convert BGR to RGB for lama
                image_rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
                result = model(image_rgb, mask, config)
                return cv2.cvtColor(result, cv2.COLOR_RGB2BGR)
            except ImportError:
                # Fall back to OpenCV method
                method = "telea"

        if method == "telea":
            return cv2.inpaint(image, mask, inpaint_radius, cv2.INPAINT_TELEA)
        else:  # ns (Navier-Stokes)
            return cv2.inpaint(image, mask, inpaint_radius, cv2.INPAINT_NS)

    def _estimate_font_color(
        self,
        image_path: str,
        region: dict[str, Any],
    ) -> tuple[int, int, int]:
        """Estimate the font color from a text region.

        Args:
            image_path: Path to the image
            region: Text region dict with box coordinates

        Returns:
            RGB tuple of the estimated font color
        """
        import cv2
        import numpy as np

        image = cv2.imread(image_path)
        box = region["box"]

        # Extract the text region
        x, y, w, h = box["x"], box["y"], box["width"], box["height"]
        roi = image[y:y+h, x:x+w]

        if roi.size == 0:
            return (0, 0, 0)  # Default to black

        # Convert to grayscale to find text pixels (usually darker or lighter than background)
        gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)

        # Use Otsu's threshold to separate text from background
        _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)

        # Determine if text is dark on light or light on dark
        mean_val = np.mean(gray)
        if mean_val > 127:
            # Light background, text is in dark pixels
            text_mask = binary == 0
        else:
            # Dark background, text is in light pixels
            text_mask = binary == 255

        # Get average color of text pixels
        if np.any(text_mask):
            text_pixels = roi[text_mask]
            avg_color = np.mean(text_pixels, axis=0)
            # Convert BGR to RGB
            return (int(avg_color[2]), int(avg_color[1]), int(avg_color[0]))

        return (0, 0, 0)  # Default to black

    @tool()
    async def images_to_editable_pptx(
        self,
        image_paths: list[str],
        output_path: str,
        lang: str = "ch",
        inpaint_method: str = "telea",
        slide_width_inches: float = 13.333,
        slide_height_inches: float = 7.5,
        min_font_size: int = 10,
        max_font_size: int = 44,
    ) -> dict[str, Any]:
        """Convert slide images to an editable PowerPoint presentation.

        This tool:
        1. Uses OCR to detect text and its positions in each image
        2. Removes text from images using inpainting to create clean backgrounds
        3. Generates a PPTX with background images and editable text boxes

        Args:
            image_paths: List of paths to slide images (in order)
            output_path: Path for the output .pptx file
            lang: OCR language ('ch' for Chinese+English, 'en' for English only)
            inpaint_method: Background reconstruction method ('telea', 'ns', or 'lama')
            slide_width_inches: Slide width in inches (default: 16:9 widescreen)
            slide_height_inches: Slide height in inches
            min_font_size: Minimum font size in points
            max_font_size: Maximum font size in points

        Returns:
            Status with output path and processing details
        """
        try:
            import cv2
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError as e:
            missing = str(e).split("'")[1] if "'" in str(e) else str(e)
            return {
                "success": False,
                "error": f"Missing dependency: {missing}",
                "hint": "Install with: pip install python-pptx opencv-python paddlepaddle paddleocr",
            }

        output_file = Path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)

        # Create presentation with specified dimensions
        prs = Presentation()
        prs.slide_width = Inches(slide_width_inches)
        prs.slide_height = Inches(slide_height_inches)

        # Use blank layout
        blank_layout = prs.slide_layouts[6]  # Blank slide

        processed_slides = []
        temp_files = []  # Track temp files for cleanup

        try:
            for idx, image_path in enumerate(image_paths):
                slide_info = {
                    "page": idx + 1,
                    "image": image_path,
                    "text_regions": 0,
                    "status": "processing",
                }

                if not Path(image_path).exists():
                    slide_info["status"] = "error"
                    slide_info["error"] = "Image file not found"
                    processed_slides.append(slide_info)
                    continue

                # Step 1: Detect text regions
                try:
                    text_regions = self._detect_text_regions(image_path, lang)
                    slide_info["text_regions"] = len(text_regions)
                except Exception as e:
                    slide_info["status"] = "error"
                    slide_info["error"] = f"OCR failed: {str(e)}"
                    processed_slides.append(slide_info)
                    continue

                # Step 2: Create background image (with text removed)
                image = cv2.imread(image_path)
                img_height, img_width = image.shape[:2]

                if text_regions:
                    # Create mask and inpaint
                    mask = self._create_text_mask(
                        (img_width, img_height),
                        text_regions,
                        dilate_pixels=8,
                    )
                    background = self._inpaint_background(
                        image_path, mask, method=inpaint_method
                    )
                else:
                    # No text detected, use original image
                    background = image

                # Save background to temp file
                temp_bg_path = output_file.parent / f"_temp_bg_{idx}.png"
                cv2.imwrite(str(temp_bg_path), background)
                temp_files.append(temp_bg_path)

                # Step 3: Create slide with background
                slide = prs.slides.add_slide(blank_layout)

                # Add background image (full slide size)
                slide.shapes.add_picture(
                    str(temp_bg_path),
                    Inches(0),
                    Inches(0),
                    width=prs.slide_width,
                    height=prs.slide_height,
                )

                # Step 4: Add text boxes for each detected text region
                # Calculate scale factors (image pixels to slide inches)
                scale_x = slide_width_inches / img_width
                scale_y = slide_height_inches / img_height

                for region in text_regions:
                    box = region["box"]

                    # Convert pixel coordinates to inches
                    left = Inches(box["x"] * scale_x)
                    top = Inches(box["y"] * scale_y)
                    width = Inches(box["width"] * scale_x)
                    height = Inches(box["height"] * scale_y)

                    # Add text box
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    tf = textbox.text_frame
                    tf.word_wrap = False

                    # Add text
                    p = tf.paragraphs[0]
                    p.text = region["text"]

                    # Set font size (scaled and clamped)
                    font_size = int(region["estimated_font_size"] * scale_y * 72)
                    font_size = max(min_font_size, min(font_size, max_font_size))
                    p.font.size = Pt(font_size)

                    # Estimate and set font color
                    try:
                        r, g, b = self._estimate_font_color(image_path, region)
                        p.font.color.rgb = RGBColor(r, g, b)
                    except Exception:
                        # Default to black on error
                        p.font.color.rgb = RGBColor(0, 0, 0)

                    # Center align (common for slides)
                    p.alignment = PP_ALIGN.LEFT

                slide_info["status"] = "success"
                processed_slides.append(slide_info)

            # Save presentation
            prs.save(str(output_file))

            # Cleanup temp files
            for temp_file in temp_files:
                try:
                    temp_file.unlink()
                except Exception:
                    pass

            # Summary
            success_count = sum(1 for s in processed_slides if s["status"] == "success")
            total_text_regions = sum(s.get("text_regions", 0) for s in processed_slides)

            return {
                "success": True,
                "message": f"Created editable PPTX with {success_count}/{len(image_paths)} slides",
                "output_path": str(output_file),
                "total_slides": len(image_paths),
                "successful_slides": success_count,
                "total_text_regions": total_text_regions,
                "slides": processed_slides,
            }

        except Exception as e:
            # Cleanup temp files on error
            for temp_file in temp_files:
                try:
                    temp_file.unlink()
                except Exception:
                    pass

            return {
                "success": False,
                "error": f"PPTX generation failed: {str(e)}",
            }

    @tool()
    async def notebookllm_close_browser(self) -> dict[str, Any]:
        """Close the browser connection (does not close the actual browser).

        Returns:
            Status of the operation
        """
        try:
            await self._close_browser()
            return {
                "success": True,
                "message": "Browser connection closed successfully",
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
            }
